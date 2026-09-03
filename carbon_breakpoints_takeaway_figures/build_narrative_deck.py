#!/usr/bin/env python3
"""Paris at Ten — narrative commentary version.
Reframes the Paris Agreement carbon intensity analysis as a three-act narrative."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

FIGS = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/outputs/final_master/figures_world_v3"
OUT = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/paris_at_ten_narrative_deck.pptx"

INK = RGBColor(0x21, 0x21, 0x21); GRAY = RGBColor(0x55, 0x55, 0x55); LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x1B, 0x78, 0x37); BLUE = RGBColor(0x3A, 0x66, 0xA5)
RED = RGBColor(0xB3, 0x3A, 0x3A); AMBER = RGBColor(0xC7, 0x7F, 0x00)

prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
slide = lambda: prs.slides.add_slide(BLANK)

def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         space_after=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, "Calibri"
    r.font.color.rgb = color
    return p

def runs(tf, parts, size=14, first=False, space_after=6, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    for t, c, b in parts:
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, "Calibri"
        r.font.italic = italic; r.font.color.rgb = c
    return p

def section_title(s, act_num, act_name, subtitle):
    tf = box(s, 0.55, 2.5, 12.25, 2.5)
    para(tf, act_num, size=72, bold=True, color=GREEN, first=True, space_after=0)
    para(tf, act_name, size=48, bold=True, color=INK, space_after=8)
    para(tf, subtitle, size=20, color=GRAY, italic=True, space_after=0)

def fit_image(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih); w, h = iw * sc, ih * sc
    return s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2), Inches(w), Inches(h))

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

F = lambda n: os.path.join(FIGS, n)

# OPENING SLIDE
s = slide()
tf = box(s, 0.9, 1.6, 11.5, 3.2)
para(tf, "Paris at Ten", size=54, bold=True, first=True, space_after=6)
para(tf, "A world aggregate study asking whether the Paris Agreement can claim credit for the recent acceleration in carbon intensity",
     size=18, color=GRAY, space_after=12)
para(tf, "World-level series, 1965–2023 · 44-specification robustness grid · unknown-break-date testing",
     size=14, color=GREEN, space_after=0)
notes(s, "The Paris Agreement was signed in 2015 as a landmark commitment to limit global warming. But climate action is not the only force that shapes emissions. This paper asks: does the timing of carbon intensity change support the treaty's attribution, or do other factors better explain what happened?")

# THE QUESTION SLIDE
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.2)
para(tf, "THE QUESTION", size=20, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.8, 12.25, 5.4)
para(tf, "Development policies seek to reverse emissions growth. But attribution is hard. Global CO2 has a dozen moving parts — economic growth, technology, policy, accident. Paris claimed all of them in 2015. The question is whether the data agrees.",
     size=16, color=GRAY, first=True, space_after=14)
para(tf, "A real turn happened: carbon intensity fell from +0.4%/yr in the 2000s to -2.1%/yr since 2012. That improvement is large, genuine, and necessary. The narrow question we ask is whether 2015 marks when the curve bent, or whether the curve had already started bending before Paris was written.",
     size=16, color=GRAY, space_after=14)
para(tf, "If Paris was the cause, the timing question is the first place where that claim should show up in the data.",
     size=16, color=INK, bold=True, space_after=0)
notes(s, "The setup. Paris in 2015 is a treaty that aims to govern emissions. But the acceleration in decarbonization began earlier — in 2012. Knowing when the turn happened is the first step to knowing what caused it.")

# ACT I - INTRODUCTION
section_title(s := slide(), "I", "ACT I", "Testing whether the timing is robust")
notes(s, "In Act I we show that the 2012 turning point is not an artifact of how the data is built. It appears in 44 different specifications. That robustness matters because it shifts the burden of proof: Paris has to explain not just why 2015 looks like a break, but why almost every reasonable construction of the data puts the break three years earlier.")

# ACT I - THE LONG RECORD
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "I. REPLICATION", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 6.0, 5.7)
para(tf, "The decisive turn is the 1970s, not a treaty", size=18, bold=True, color=INK, first=True, space_after=10)
para(tf, "The world's carbon intensity record shows not one turn but three. The largest by far came in the 1970s with the oil shocks, cutting the growth rate from +0.48%/yr to -1.54%/yr. Another came around 2000, a smaller one. The current fast episode started in 2012, at -2.11%/yr, faster than any prior episode but part of a longer pattern.",
     size=14, color=GRAY, space_after=12)
para(tf, "The 2012 turn is statistically significant (p = 0.002) and robust: 89% of our different ways of building the data place it before 2015. What does not support 2015 is the year that matters for Paris attribution.",
     size=14, color=GRAY, space_after=0)
fit_image(s := s, FIGS + "/v3_fig1_long_record.png", 6.8, 1.55, 6.3, 5.7)
notes(s, "Figure 1 shows the whole record. The 1970s oil shocks were the largest break; the 2012 turn is faster but third in order of magnitude. The formal unknown-break test puts the dominant break at 1972 for overall intensity (p=0.002) and 1974 for efficiency (p<0.001).")

# ACT I - THE ROBUSTNESS TEST
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "I. THE ROBUSTNESS TEST", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 6.5, 5.7)
para(tf, "Across 44 specifications, the turn comes before Paris", size=18, bold=True, color=INK, first=True, space_after=10)
para(tf, "We ran the same analysis 44 times: different starting years, endpoints, how much data we require, whether we fill gaps in the data. Different researchers might make these choices differently. What we found is that 89% of those 44 runs place the best-fitting turning point before 2015, usually between 2011 and 2013.",
     size=14, color=GRAY, space_after=12)
para(tf, "This is the single most robust result in the analysis. It does not matter much how we build the world aggregate — the timing question is not a choice about data.",
     size=14, color=GRAY, space_after=0)
fit_image(s := s, FIGS + "/v3_fig3_specification_invariance.png", 6.85, 1.55, 6.15, 5.7)
notes(s, "Each dot is one specification. Almost all cluster before 2015. The red line at 2015 shows Paris; the blue line at 1992 shows Rio. The timing of the turn is visible, consistent, and not at Paris.")

# WHERE ACT I LEAVES US
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.2)
para(tf, "WHERE ACT I LEAVES US", size=20, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.8, 12.25, 5.4)
para(tf, "The robustness of the 2012 timing means the attribution question is now specific. Paris does not have to explain why one particular construction of the data looks significant at 2015. It has to explain why almost every reasonable construction of the data places the turn three years before Paris was written.",
     size=16, color=GRAY, first=True, space_after=14)
para(tf, "That shifts the evidentiary standard. A timing coincidence is less impressive when the timing is robust and inconvenient. We now move to Act II to ask a second question: what mechanism changed in 2012? Was it the fuel mix, energy efficiency, or something about growth? Understanding what changed will clarify what could have changed it.",
     size=16, color=GRAY, space_after=0)
notes(s, "Act I established that timing is not an artifact. Now we move to mechanism. The Kaya identity tells us there are only three moving parts to carbon intensity — growth, fuel mix, and efficiency. We'll examine each one.")

# ACT II - INTRODUCTION
section_title(s := slide(), "II", "ACT II", "Understanding what changed: the fuel mix stalled, then resumed")
notes(s, "In Act II we use the Kaya identity — the decomposition of carbon intensity into fuel mix and efficiency — to understand the mechanism. This matters because it tells us where policy levers operate and where they do not.")

# ACT II - THE DECOMPOSITION
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "II. THE MECHANISM", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 6.2, 5.7)
para(tf, "Efficiency stayed flat. The fuel mix stalled, then resumed.", size=18, bold=True, color=INK, first=True, space_after=10)
para(tf, "Carbon intensity has two doors: the fuel mix (how much carbon per unit of energy) and efficiency (how much energy per unit of GDP). Efficiency improved steadily at about -1% per year for the entire fifty-year period — unchanged by any event, including Paris.",
     size=14, color=GRAY, space_after=12)
para(tf, "The fuel mix tells a different story. It improved at -0.52%/yr after the oil shocks, stalled completely from 1990 to 2015, then resumed at -0.70%/yr. The resumption looks like a return to historical pace, not a new policy effect.",
     size=14, color=GRAY, space_after=0)
fit_image(s := s, FIGS + "/v3_fig2_baseline_problem.png", 6.5, 1.55, 6.5, 5.7)
notes(s, "The three eras are 1973-1990 (oil shock recovery), 1990-2015 (Rio to Paris stall), 2016-2023 (resumption). Efficiency steady throughout. Fuel mix is the story.")

# ACT II - WHAT THIS MEANS
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "II. THE DIAGNOSIS", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 12.25, 5.7)
para(tf, "Why the fuel mix matters for attribution", size=18, bold=True, color=INK, first=True, space_after=10)
para(tf, "The mechanism is where policy lives. Climate agreements try to change the fuel mix — through renewable targets, carbon pricing, technology support. They do not directly move efficiency. If Paris was the cause of the 2012 turn, efficiency should have accelerated. It did not.",
     size=14, color=GRAY, space_after=12)
para(tf, "The fuel mix resumed in 2012 after a 25-year stall, but the resumption matches the timing of renewable cost collapse, not the treaty timeline. Solar photovoltaic costs fell by 90% between 2010 and 2020, with the steepest drops between 2010 and 2015 — exactly when the fuel mix started improving again.",
     size=14, color=GRAY, space_after=12)
para(tf, "This is the puzzle Paris has to solve: the acceleration happened where it should have happened (fuel mix) but when it should not have happened (before the treaty), and it looks like what technology-driven change produces.",
     size=14, color=GRAY, space_after=0)
notes(s, "The timing and mechanism both point to technology. The Paris timing is harder to support with the data.")

# WHERE ACT II LEAVES US
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.2)
para(tf, "WHERE ACT II LEAVES US", size=20, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.8, 12.25, 5.4)
para(tf, "The 2012 turn in the fuel mix is the core acceleration the world needs. It is a real, large, and rapid shift. But it precedes Paris, lacks an efficiency partner, and coincides with technology cost collapse. The mechanism does not point toward Paris as the cause. It points toward renewable technology getting cheap enough that switching fuel became profitable.",
     size=16, color=GRAY, first=True, space_after=14)
para(tf, "Act III now turns to the harder question: what does this mean for how we think about Paris? Not whether Paris is worthless, but what the data allows us to claim about it.",
     size=16, color=GRAY, space_after=0)
notes(s, "Now we move to interpretation. The facts are clear. The implications are harder.")

# ACT III - INTRODUCTION
section_title(s := slide(), "III", "ACT III", "What the data can and cannot support")
notes(s, "Act III draws out the implications. We do not argue that Paris is worthless. We argue that the timing question is the wrong one to use for attribution, or at least the available data does not support that reading.")

# ACT III - WHAT WE CAN SAY
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "III. THE CONCLUSION", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 6.0, 5.7)
para(tf, "What the data supports", size=18, bold=True, color=GREEN, first=True, space_after=10)
para(tf, "The world is decarbonizing faster than ever before. Carbon intensity is declining at -2.1%/yr since 2012, faster than any prior episode in the measured record. This improvement is large, genuine, and necessary.",
     size=14, color=GRAY, space_after=12)
para(tf, "It is also concentrated where policy can work: the fuel mix. Efficiency contributions remain steady and small. The world has not yet responded on the demand side of emissions. That matters for whether the current pace can accelerate.",
     size=14, color=GRAY, space_after=0)

tf = box(s, 6.3, 1.55, 6.3, 5.7)
para(tf, "What the data does not support", size=18, bold=True, color=RED, first=True, space_after=10)
para(tf, "The Paris agreement was signed in 2015. The acceleration in carbon intensity began in 2012, in roughly 9 of every 10 specifications. Energy efficiency did not respond to the treaty at all. The timing does not point toward Paris as the cause.",
     size=14, color=GRAY, space_after=12)
para(tf, "This is not an argument that Paris is ineffective or unnecessary. It is an argument that the timing data do not support attribution to Paris. Other factors — technology cost, market dynamics, pre-existing policy in major economies — appear more directly responsible for when the turn happened.",
     size=14, color=GRAY, space_after=0)
notes(s, "The honest reading. Not an indictment but a dating result that does not support Paris attribution based on timing alone.")

# ACT III - WHAT SETTLES IT
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "III. THE OPEN QUESTIONS", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 12.25, 5.7)
para(tf, "What would change this conclusion?", size=18, bold=True, color=INK, first=True, space_after=12)
para(tf, "Three questions can still be answered to test whether Paris played a role. First: does the acceleration persist? A technology shock fades as solar saturates; policy ratcheting persists. Another decade will tell. Second: does efficiency ever accelerate? Technology alone cannot explain steady efficiency gains; if efficiency moves, policy becomes credible. Third: what do countries show? Paris works through national pledges. If the treaty is real, its signature should appear when you compare country by country.",
     size=14, color=GRAY, space_after=12)
para(tf, "Of those three, the country-level question is the most tractable and the most direct. It does not rely on waiting for more data or a mechanism that may never materialize. If Paris changed country behavior, the country numbers should show it.",
     size=14, color=GRAY, space_after=0)
notes(s, "These are not criticisms of Paris. They are genuinely open empirical questions about what the data can settle.")

# WHERE ACT III LEAVES US
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.2)
para(tf, "WHERE THIS LEAVES IT", size=20, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.8, 12.25, 5.4)
para(tf, "Global carbon intensity is accelerating in the right direction, at the largest scale we have ever measured. Whether Paris deserves credit depends on what we look for and how we look for it. The timing question points elsewhere — toward renewable technology and the 2010s cost collapse. The country question remains open, and it is the most direct one Paris attribution has to answer.",
     size=16, color=GRAY, first=True, space_after=14)
para(tf, "The acceleration is real and necessary. It started before Paris, in a direction that makes sense for what happened to solar and wind cost. Whether and to what degree Paris drove it remains an open question best answered at the national level.",
     size=16, color=GRAY, space_after=0)
notes(s, "Honest conclusion. Not a verdict on Paris but a parsing of what the world-level data can and cannot claim.")

# POLICY IMPLICATIONS 1
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "IMPLICATION 1 · TECHNOLOGY IS THE LEVER THAT MOVED", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 12.25, 5.7)
para(tf, "The data tells a story about what actually changes decarbonization", size=18, bold=True, color=INK, first=True, space_after=12)
para(tf, "The fuel mix resumed its decline in 2012, exactly when solar photovoltaic costs began their steepest fall. The timing is too precise to be coincidence. Renewable technology became cheap enough that switching was profitable, independent of treaty commitments.",
     size=14, color=GRAY, space_after=12)
para(tf, "This does not mean Paris is irrelevant. But it suggests that climate policy works best when it either makes clean technology cheaper (subsidies, R&D) or makes dirty energy more expensive (carbon pricing). Commitments alone, without the economic incentive to change, do not appear sufficient to move large-scale behavior.",
     size=14, color=GRAY, space_after=12)
para(tf, "For policymakers: the Copenhagen Accord and Paris Agreement happened while the solar industry was already collapsing in cost. The timing alignment is favorable but deceptive — the mechanism was technology, not treaty enforcement.",
     size=14, color=GRAY, space_after=0)
notes(s, "The policy insight: we credit Paris with decarbonization that happened because of technology cost collapse. That's not a criticism of Paris, but it means our theories of change matter. Technology is what moved behavior at world scale.")

# POLICY IMPLICATIONS 2
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "IMPLICATION 2 · EFFICIENCY IS THE MISSING HALF", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 12.25, 5.7)
para(tf, "Where climate policy has not yet delivered", size=18, bold=True, color=INK, first=True, space_after=12)
para(tf, "Energy efficiency — how much energy the world uses per unit of economic output — has improved at a steady -1% per year for fifty years, completely unmoved by Paris or any other policy intervention. This is the constraint on how fast decarbonization can go.",
     size=14, color=GRAY, space_after=12)
para(tf, "The fuel mix can only get so clean before renewable penetration hits physical or storage limits. Beyond that, the world has to use less energy per unit of output. Current evidence suggests that policy has not yet moved that needle. This is the frontier for climate strategy: understanding why efficiency is sticky and what would accelerate it.",
     size=14, color=GRAY, space_after=12)
para(tf, "For policymakers: the current world decarbonization rate of -2.1%/yr needs to reach -5% to -11%/yr to meet climate targets. The fuel mix alone cannot carry that load. Efficiency has to move, and so far it has not responded to policy.",
     size=14, color=GRAY, space_after=0)
notes(s, "Efficiency is where we haven't made progress. It's also where the big gains are available if policy can figure out how to unlock them.")

# POLICY IMPLICATIONS 3
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.0)
para(tf, "IMPLICATION 3 · ATTRIBUTION REQUIRES COUNTRY-LEVEL WORK", size=16, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.55, 12.25, 5.7)
para(tf, "Why global timing does not settle the Paris question", size=18, bold=True, color=INK, first=True, space_after=12)
para(tf, "At world level, timing is ambiguous. The global acceleration could reflect renewable cost, pre-existing climate policies in major economies (California's standards, China's targets predating Paris), or Paris itself. We cannot tell from the world aggregate alone.",
     size=14, color=GRAY, space_after=12)
para(tf, "Countries offer a cleaner test. If Paris changed national behavior, each country's carbon intensity should show a break near 2015, independent of technology cost. If the break appears everywhere before 2015, or if it does not correlate with Paris commitments, then technology cost is the simpler explanation.",
     size=14, color=GRAY, space_after=12)
para(tf, "For policymakers and researchers: do not settle the Paris question at world level. The answer lives in country-by-country variation, where you can control for technology diffusion and test whether national commitments shaped national outcomes.",
     size=14, color=GRAY, space_after=0)
notes(s, "This is why the companion country paper matters. The world-level question is genuinely ambiguous. Countries can resolve it.")

# FINAL SYNTHESIS
s = slide()
tf = box(s, 0.55, 0.5, 12.25, 1.2)
para(tf, "THE BOTTOM LINE", size=20, bold=True, color=GREEN, first=True, space_after=0)
tf = box(s, 0.55, 1.8, 12.25, 5.4)
para(tf, "The world is decarbonizing faster than ever before. That is real, necessary, and visible in the data. But the cause is not what the treaty timeline suggests. At world level, renewable technology cost collapse provides a simpler and more direct explanation for when and how the acceleration began.",
     size=16, color=GRAY, first=True, space_after=14)
para(tf, "This is not an argument that Paris is ineffective or unimportant. It is an argument that world-level attribution is insufficient to settle the question. Whether and to what degree Paris shifted national behavior requires country-level analysis where policy effects can be isolated from technology diffusion. That question is open and important. It is also answerable.",
     size=16, color=GRAY, space_after=0)
notes(s, "Final framing: we have a real acceleration, a plausible alternative explanation, and an open empirical question. That is the honest statement.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
