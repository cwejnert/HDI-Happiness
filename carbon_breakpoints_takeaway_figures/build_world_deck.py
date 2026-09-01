#!/usr/bin/env python3
"""Paris at Ten — world-level deck. Country analysis deferred to the follow-on paper."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

FIGS = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/outputs/final_master/figures_world"
OUT = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/paris_at_ten_world_deck.pptx"

INK = RGBColor(0x21, 0x21, 0x21); GRAY = RGBColor(0x55, 0x55, 0x55); LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x1B, 0x78, 0x37); BLUE = RGBColor(0x3A, 0x66, 0xA5)
RED = RGBColor(0xB3, 0x3A, 0x3A); AMBER = RGBColor(0xC7, 0x7F, 0x00)

prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
slide = lambda: prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         space_after=6, align=PP_ALIGN.LEFT, font="Calibri"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, font
    r.font.color.rgb = color
    return p


def runs(tf, parts, size=14, first=False, space_after=6, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    for t, c, b in parts:
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, "Calibri"
        r.font.italic = italic; r.font.color.rgb = c
    return p


def kicker_title(s, kicker, title, kcolor=GREEN):
    tf = box(s, 0.55, 0.3, 12.25, 1.0)
    para(tf, kicker, size=12, color=kcolor, bold=True, first=True, space_after=2)
    para(tf, title, size=24, color=INK, bold=True, space_after=0)


def fit_image(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih); w, h = iw * sc, ih * sc
    return s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2), Inches(w), Inches(h))


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


def figure_slide(kicker, title, img, narrative=None, guide=None, note="", kcolor=GREEN):
    s = slide(); kicker_title(s, kicker, title, kcolor=kcolor); top = 1.30
    if narrative:
        tf = box(s, 0.55, top, 12.25, 0.80); para(tf, narrative, size=13, color=GRAY, first=True, space_after=0)
        top += 0.84
    if guide:
        tf = box(s, 0.55, top, 12.25, 0.5)
        runs(tf, [("How to read it: ", BLUE, True), (guide, LGRAY, False)], size=11.5, first=True, space_after=0, italic=True)
        top += 0.52
    fit_image(s, img, 0.35, top + 0.04, 12.63, 7.2 - top)
    notes(s, note); return s


def act_open(roman, label, title, lede, stats, note=""):
    s = slide()
    tf = box(s, 0.9, 0.85, 11.5, 2.6)
    para(tf, roman, size=44, bold=True, color=GREEN, first=True, space_after=0)
    para(tf, label, size=12, bold=True, color=LGRAY, space_after=6)
    para(tf, title, size=28, bold=True, space_after=8)
    para(tf, lede, size=14, color=GRAY, space_after=0)
    x = 0.9
    for big, small in stats:
        tf = box(s, x, 4.7, 3.8, 2.0)
        para(tf, big, size=28, bold=True, color=GREEN, first=True, space_after=2)
        para(tf, small, size=12.5, color=GRAY, space_after=0)
        x += 3.95
    notes(s, note); return s


F = lambda n: os.path.join(FIGS, n)

# 1 · title
s = slide()
tf = box(s, 0.9, 1.6, 11.5, 3.7)
para(tf, "Paris at Ten", size=46, bold=True, first=True, space_after=8)
para(tf, "The world's carbon-intensity curve is bending faster than at any time on record — and for the first time, a treaty date survives the placebo test. It still is not the best-fitting one.",
     size=18, color=GRAY, space_after=14)
runs(tf, [("Structural breaks · Kaya pathways · the global record, 1971-2021", GREEN, True)], size=14)
tf = box(s, 0.9, 5.95, 11.5, 0.9)
para(tf, "58-country composition-constant panel · 87% of global CO2 · four estimated episodes · three Kaya pathways",
     size=13, color=INK, bold=True, first=True)
notes(s, "World-level paper. The country-level mechanism analysis (97 breaks, 7 mechanisms, the nine) is deferred to the follow-on paper and referenced only at the end.")

# 2 · the setup
s = slide()
kicker_title(s, "THE SETUP", "Emissions growth has nearly stopped. Can we say Paris is why?")
tf = box(s, 0.55, 1.55, 7.35, 5.5)
para(tf, "The good news is real. Worst-case scenarios have moved off the table, and CO2 growth across the countries we can track consistently fell from +2.5% per year over the Rio-to-Paris quarter century to +0.7% per year since 2015. On a constant country set, carbon intensity is now falling at -2.2% a year: the fastest sustained decarbonization anywhere in the record.",
     size=14, color=GRAY, first=True, space_after=12)
para(tf, "But there is an obvious way to over-read that success — to attribute it to the agreement whose anniversary it coincides with. Carbon intensity has been falling for fifty years, in episodes with their own causes. So the question is narrow and awkward: can the global record distinguish a Paris effect from an acceleration already under way when the treaty was signed?",
     size=14, color=GRAY, space_after=12)
para(tf, "We do not assume treaty years are turning points. We let the data choose the year, then ask how much better that year is than its neighbours — and hold Paris to the standard we would hold any other date.",
     size=14, color=INK, space_after=0)
tf = box(s, 8.25, 1.55, 4.55, 5.5)
para(tf, "Why intensity, not emissions", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "Emissions arithmetic has three moving parts: economic growth, the fuel mix, and energy per unit of output. Growth pushes emissions up; the two intensity terms are the only levers that push them down.",
     size=13, color=GRAY, space_after=10)
para(tf, "So intensity is the leading indicator. A policy signal must appear in the intensity slopes before it can appear in emissions at all — and unless intensity improves faster than the economy grows, emissions cannot stabilize.",
     size=13, color=INK, space_after=0)
notes(s, "Sympathetic frame, then the identity that justifies the design. Note the numbers here are balanced-panel: +2.5 to +0.7 CO2 growth, -2.2%/yr intensity in the current episode.")

# 3 · the argument
s = slide()
kicker_title(s, "THE ARGUMENT", "Four acts")
acts = [("I", "What the record looks like", "Four episodes, not one turning point — and the current one is the fastest. But the world series has to be built carefully before any of it means anything."),
        ("II", "The anniversary test", "Since 2015 intensity has run below its own extrapolated trend, significantly, through both Kaya doors."),
        ("III", "Can the record name the author?", "Rio the record dismisses. Paris it cannot — yet the best-fitting onset is still two years early, and the test cannot resolve two years."),
        ("IV", "What would settle it", "The evidence that would decide this by Paris+20 — and why it has to be gathered country by country.")]
x = 0.55
for roman, t, d in acts:
    tf = box(s, x, 1.7, 2.95, 4.8)
    para(tf, roman, size=38, bold=True, color=GREEN, first=True, space_after=6)
    para(tf, t, size=15.5, bold=True, space_after=8)
    para(tf, d, size=12.5, color=GRAY, space_after=0)
    x += 3.12
notes(s, "Act III is the pivot: the finding is not a null, it is an unresolvable-but-favourable coincidence. Act IV hands off to the national paper.")

# 4 · the instrument
s = slide()
kicker_title(s, "THE INSTRUMENT", "Two straight lines and a kink", kcolor=BLUE)
tf = box(s, 0.55, 1.55, 6.1, 5.5)
para(tf, "Every series is plotted in logs, so a straight line means a steady percentage change per year. We fit two straight lines that meet at a kink, trying every admissible year and keeping the one that fits best. That year is the data's own estimate of when the trend changed, chosen with no reference to any treaty.",
     size=13.5, color=GRAY, first=True, space_after=12)
para(tf, "The date is then stress-tested — re-estimated on reshuffled versions of the series — and, crucially, compared against every nearby year. A break whose year barely moves and clearly beats its neighbours is credible; one that sits on a plateau of equally good years is not, however significant it looks.",
     size=13.5, color=GRAY, space_after=12)
para(tf, "(Formally: y(t) = a + b·t + d·max(t−k, 0). Autocorrelation-robust standard errors; BIC for model evidence and for how many breaks the record supports.)",
     size=10.5, color=LGRAY, italic=True, space_after=0)
tf = box(s, 7.0, 1.55, 5.8, 5.5)
para(tf, "Two doors for any improvement", size=15, bold=True, first=True, space_after=6)
para(tf, "Carbon intensity — CO2 per unit of output — can only improve through one of two doors: cleaner energy (less CO2 per unit of energy: fuel switching, nuclear, renewables, the power system), or less energy per unit of output (efficiency, and changes in what the economy makes).",
     size=13.5, color=GRAY, space_after=10)
para(tf, "We estimate the bend in the overall series and behind both doors. Which door moves is diagnostic: a fuel-mix bend points to energy policy and technology; an efficiency bend points as often to economic structure. A treaty effect should show up in the supply door first.",
     size=13.5, color=INK, space_after=0)
notes(s, "Keep it conversational. The 'two doors' framing pays off in Act II, where both doors turn out to be moving.")

# 5 · rules of evidence
s = slide()
kicker_title(s, "THE RULES OF EVIDENCE", "What it takes to credit a treaty — and to trust a world series", kcolor=BLUE)
tf = box(s, 0.55, 1.55, 5.9, 5.5)
para(tf, "When a treaty year earns credit", size=15, bold=True, first=True, space_after=6)
para(tf, "The catch with treaty-year tests is that a gently curving trend makes almost any year look significant. So a treaty year earns credit only if the slope genuinely changes there by a robust test, the treaty year beats nearly all nearby placebo years, the best-fitting date lands within a few years of it, and there is a plausible mechanism.",
     size=13, color=GRAY, space_after=10)
para(tf, "Rio fails the second test badly. Paris passes it. That difference is the paper.", size=13, color=INK, bold=True, space_after=0)
tf = box(s, 6.85, 1.55, 5.95, 5.5)
para(tf, "And what counts as 'the world'", size=15, bold=True, first=True, space_after=6)
para(tf, "A world series summed over whichever countries reported that year is not a world series — it mixes trend with coverage. In this database coverage jumps from 59 to 149 countries in 1980 and collapses from 191 to 78 in 2022 as reporting lags.",
     size=13, color=GRAY, space_after=10)
para(tf, "Everything here therefore uses a composition-constant panel: the 58 countries with complete data every year from 1971 to 2021, 87% of global CO2. This is not a technicality — it doubles the estimated post-Paris acceleration and turns the demand side from insignificant to significant.",
     size=13, color=INK, space_after=0)
notes(s, "The balanced-panel point is a genuine methods contribution for a world-level paper: naive aggregation understated the recent acceleration.")

# 6 · ACT I
act_open("I", "ACT I", "What the record actually looks like",
         "Before asking whether Paris bent the curve, establish what the curve has been doing — on a country set that does not change underneath the question.",
         [("4 episodes", "the record prefers three breaks (1984, 2002, 2012) over one"),
          ("-0.07%/yr", "carbon intensity through the 2002-2012 plateau — decarbonization stalled"),
          ("-2.20%/yr", "since 2012: the fastest sustained improvement in the record")],
         note="The stall-then-surge shape is the single most useful thing this paper can put on the record.")

figure_slide("ACT I · THE EPISODES", "Decarbonization stalled for a decade, then resumed faster than ever",
             F("world_fig01_episodes.png"),
             narrative="The record is not one turning point but a sequence: steady improvement to 1984, slower improvement to 2002, a decade-long plateau while coal-fired industrial growth offset efficiency gains everywhere else, and then the sharpest sustained decarbonization in the series. Paris arrives three years into that fourth episode.",
             guide="Grey dots are the composition-constant panel's carbon intensity (1971 = 100, log scale); the red line is the best multiple-break fit, with each segment's average annual change labelled.",
             note="BIC prefers three breaks over one. The 2002-2012 plateau is the China/WTO coal era; the post-2012 episode is the renewables cost collapse plus post-crisis restructuring.")

figure_slide("ACT I · THE AGGREGATION TRAP", "A world series summed over changing countries is not a world series",
             F("world_fig05_aggregation.png"),
             narrative="Coverage in this database is not constant — it jumps by ninety countries in 1980 and collapses by more than half in 2022 as reporting lags. Summing whatever is available each year mixes real trend with bookkeeping. And even on a fixed panel, the world average conceals its largest member: China's coal-era growth held the global fuel mix flat for two decades while everyone else slowly improved.",
             guide="Left: countries with complete Kaya data each year. Right: carbon intensity of energy for the fixed panel, with and without China.",
             note="This is why the naive aggregate understated the post-Paris acceleration: the 2022-23 points were computed on 78 countries, and the pre-2015 baseline was flattened by China.",
             kcolor=BLUE)

# 9 · ACT II
act_open("II", "ACT II", "The anniversary test",
         "Extend the Rio-to-Paris trend past 2015 and ask whether the world has beaten it.",
         [("-1.86 pp/yr", "faster than trend since Paris, overall carbon intensity (t = 4.9)"),
          ("both doors", "fuel mix and energy intensity now accelerate significantly"),
          ("about half", "of the emissions slowdown is intensity, not slower growth")],
         note="On the naive aggregate the demand side looked flat; on the balanced panel it is significant. That correction matters.")

figure_slide("ACT II · THE ANNIVERSARY TEST", "The world is running well ahead of its own trend",
             F("world_fig02_anniversary_test.png"),
             narrative="Every point since 2015 falls below the extrapolated Rio-to-Paris trend, in all three series, and the gap widens each year. This is not a marginal result: overall intensity is improving 1.9 percentage points a year faster than its prior trend. And it runs through both doors — the fuel mix and energy intensity are both accelerating, which is what a broad transition rather than a single technology looks like.",
             guide="1990 = 100, log scale. Solid line is the 1990-2015 trend; dashed extends it with a 95% band for where years should fall if nothing had changed; red points are observed since Paris.",
             note="C/GDP -1.86 (t=4.9), C/E -0.84 (t=4.2), E/GDP -1.02 (t=4.8). Every best-fitting onset lands before 2015 — that is Act III.")

figure_slide("ACT II · THE ARITHMETIC", "This time, intensity is doing about half the work",
             F("world_fig03_arithmetic.png"),
             narrative="Emissions growth fell by 1.9 percentage points after Paris. Roughly half of that is slower economic growth — including the pandemic — and roughly half is faster intensity improvement, split between a cleaner fuel mix and lower energy intensity. That is a materially better mix than the Rio-to-Paris era, when the fuel mix contributed nothing at all. But the dotted line shows the target: intensity must fall as fast as the economy grows before emissions stop rising.",
             guide="Each bar is an average annual growth rate; by construction the red emissions bar equals the grey GDP bar plus the two coloured intensity bars.",
             note="Rio-to-Paris: GDP +3.09, C/E +0.02, E/GDP -0.58, CO2 +2.52. After Paris: GDP +2.14, C/E -0.49, E/GDP -0.99, CO2 +0.65.")

# 12 · ACT III
act_open("III", "ACT III", "Can the record name the author?",
         "An acceleration is not an attribution. The test is whether the treaty year beats the years around it — the same test that dismissed Rio.",
         [("7th vs 97th", "percentile of candidate years: Rio versus Paris, overall intensity"),
          ("2012-2014", "where the best-fitting onset lands in every component"),
          ("~0.8 pp/yr", "the smallest effect ten years of data could have detected")],
         note="This is the paper's core. Paris is the first treaty date the global record cannot dismiss — and still cannot confirm.", )

figure_slide("ACT III · RIO VERSUS PARIS", "Rio the record dismisses. Paris it cannot.",
             F("world_fig04_placebo_plateau.png"),
             narrative="Place the break at every candidate year in turn and the two treaties look nothing alike. Rio sits at the 7th percentile of candidate years for overall intensity — which is why 1992 results never survive scrutiny. Paris sits at the 97th: the first treaty date this test does not throw out. But in none of the three components is it the best-fitting year, and what it sits on is a plateau, not a peak.",
             guide="Evidence for a slope change placed at each candidate year; open circles mark 2015. Higher means the record prefers that year as the onset.",
             note="Window-sensitive: on a 1990-start window Paris drops to the 79th percentile for C/GDP but holds at 95th for C/E. Report both.",
             kcolor=AMBER)

figure_slide("ACT III · THE POWER PROBLEM", "Ten years of data can only reveal a very large effect",
             F("world_fig06_power.png"),
             narrative="The observed acceleration clears the bar because it is unusually large. An effect half its size would not be visible until the early 2030s — so the record could not have detected a modest Paris effect had one existed, and it certainly cannot resolve a two-year difference in onset. The instrument is blunt at this horizon.",
             guide="The smallest post-2015 slope change detectable at 80% power (solid) and at the significance threshold (dashed), as more years accumulate; the red line is what we actually observe.",
             note="This is the honest reason the global record cannot settle attribution — not that Paris failed, but that the instrument is blunt at this horizon.",
             kcolor=AMBER)

# 15 · verdict
s = slide()
kicker_title(s, "WHERE ACT III LEAVES US", "The strongest defensible claim", kcolor=AMBER)
tf = box(s, 0.55, 1.6, 6.05, 5.4)
para(tf, "What the record establishes", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "The world is decarbonizing faster than at any point in the observed record, through both Kaya doors, and the improvement is large and statistically unambiguous. Paris is the first climate treaty whose date survives comparison with its neighbouring years — a test Rio fails badly. Whatever is happening began in the same window as the agreement and is entirely consistent with it.",
     size=13.5, color=GRAY, space_after=0)
tf = box(s, 7.0, 1.6, 5.8, 5.4)
para(tf, "What it cannot establish", size=15, bold=True, color=RED, first=True, space_after=6)
para(tf, "The best-fitting onset is 2012-2014, before the treaty, and the evidence forms a plateau across those years rather than a peak at 2015. With ten years of data the test cannot resolve differences of two or three years, and could not have seen a modest effect at all. The acceleration also coincides exactly with the collapse in renewable costs — a rival explanation the global series has no way to separate.",
     size=13.5, color=GRAY, space_after=8)
para(tf, "Paris is consistent with the acceleration, and indistinguishable from the technology shift it helped consolidate.",
     size=13.5, color=INK, bold=True, space_after=0)
notes(s, "If one sentence survives the talk, it is the bolded one. Neither vindication nor indictment: an honest 'consistent with, not identified'.")

# 16 · ACT IV
act_open("IV", "ACT IV", "What would settle it",
         "The global aggregate has told us what it can. Three things would decide the question — and two of them require leaving the world level.",
         [("Paris + 20", "when a modest effect first becomes detectable in this series"),
          ("both doors", "a demand-side acceleration is harder to attribute to cheap solar"),
          ("country by country", "where placebo tests have the power the aggregate lacks")],
         note="This is the bridge to the follow-on paper without spending slides on its content.")

s = slide()
kicker_title(s, "THE ASK", "Judge the regime where its signal must appear")
tf = box(s, 0.55, 1.65, 11.9, 5.0)
para(tf, "If the argument holds, the right test of a climate agreement is not whether emissions charts bend at signing ceremonies. Emissions are the last thing to move. The test is whether intensity slopes accelerate, in the supply door first, at a date that beats its own neighbours — and that test needs to be run where it has power.",
     size=15, color=GRAY, first=True, space_after=12)
para(tf, "Three things would decide this. Time: at current noise levels, an effect half the size of what we observe becomes visible in this series around 2030, so Paris+20 is a genuine decision point rather than a rhetorical one. Pathway: the demand-side acceleration is the more informative half, because it is much harder to attribute to falling solar costs than the fuel-mix acceleration is. And disaggregation: the same placebo discipline, applied country by country, has the statistical power and the documentary record that the global aggregate lacks — which is where this work goes next.",
     size=15, color=GRAY, space_after=12)
para(tf, "The question we would put to the field: should treaty assessment be built on intensity breakpoints and placebo discipline, rather than on emissions levels and anniversaries?",
     size=15, color=INK, bold=True, space_after=0)
notes(s, "Forward-looking and methodological. Ends on the question the paper exists to pose.")

# 18 · limits
s = slide()
kicker_title(s, "THE LIMITS", "What this design supports, and what it refuses to claim", kcolor=RED)
col = [("Composition and coverage",
        "Results are reported on a 58-country constant panel covering 87% of global CO2, not on all reporting countries. That excludes most of Africa and many small states, and it ends in 2021 because 2022-23 reporting is incomplete. The choice is defensible but it is a choice, and the naive alternative gives materially weaker estimates."),
       ("Rival explanations",
        "The post-2012 acceleration coincides with the collapse in solar and battery costs, the post-crisis slowdown in Chinese heavy industry, and coal-to-gas substitution in several large economies. The global series cannot separate these from treaty effects, and we do not claim it can."),
       ("What a breakpoint is",
        "A breakpoint identifies when a trend changed, never why. Even a perfect fit at 2015 would be a coincidence of timing, not a causal estimate. And a treaty that prevented backsliding would produce no break at all — absence of a break is not absence of an effect."),
       ("Window sensitivity",
        "Paris's standing depends on the comparison window: it sits at the 97th percentile of candidate years on a 1980-start series and the 79th on a 1990-start one for overall intensity, though it holds above the 95th for the fuel mix in both. We report both rather than choosing the flattering one.")]
xs, ys = [0.55, 7.0, 0.55, 7.0], [1.55, 1.55, 4.3, 4.3]
for (t, d), x, y in zip(col, xs, ys):
    tf = box(s, x, y, 5.9, 2.65)
    para(tf, t, size=14.5, bold=True, first=True, space_after=4)
    para(tf, d, size=12, color=GRAY, space_after=0)
notes(s, "The credibility slide. The window-sensitivity admission is the one a referee will look for — better volunteered.")

# 19 · what's next
s = slide()
kicker_title(s, "WHAT'S NEXT", "The same discipline, where it has power")
tf = box(s, 0.55, 1.7, 11.9, 4.6)
para(tf, "The global aggregate is the wrong instrument for treaty attribution, and this paper's own results explain why: eight to ten post-treaty observations, an onset plateau spanning the treaty year, and rival explanations moving on the same schedule. Paris was designed to work through national pledges and national transformation, so its signature should appear country by country — where the series are longer, the placebo tests are independent, and a documentary record exists to say what accompanied each bend.",
     size=15, color=GRAY, first=True, space_after=12)
para(tf, "That is the companion paper: the same breakpoint and placebo machinery applied to every country with adequate data, with historically coded mechanisms and a transition-quality framework that separates genuine decarbonization from recession, disruption, and offshoring.",
     size=15, color=GRAY, space_after=12)
runs(tf, [("This paper establishes the puzzle at the global level. The national paper is where it gets resolved.", INK, True)], size=15)
notes(s, "One slide only. Tees up paper 2 without previewing its results — the 97 breaks, seven mechanisms, and nine constructive cases live there.")

# 20 · close
s = slide()
tf = box(s, 0.9, 1.5, 11.5, 2.0)
para(tf, "The curve is bending faster than ever.", size=27, bold=True, first=True, space_after=2)
para(tf, "Paris is the first treaty the record cannot dismiss — or confirm.", size=27, bold=True, color=GREEN, space_after=0)
tf = box(s, 0.9, 3.6, 11.5, 2.9)
for t in ["Decarbonization stalled through 2002-2012 and then resumed at -2.2% a year, the fastest sustained improvement in the record.",
          "Since 2015 the world has run significantly ahead of its own trend, through both the fuel mix and energy intensity.",
          "Rio sits at the 7th percentile of candidate onset years; Paris sits at the 97th — but the best fit is 2012-2014, and ten years of data cannot resolve two.",
          "Consistent with, and not separable from, the technology shift it helped consolidate — which is why the next test is national."]:
    runs(tf, [("—  ", GREEN, True), (t, GRAY, False)], size=14, space_after=9)
tf = box(s, 0.9, 6.55, 11.5, 0.7)
para(tf, "Structural breaks say when · Kaya says through which door · placebo discipline says whether the date is special · power says what we could have seen",
     size=13, italic=True, color=LGRAY, first=True, align=PP_ALIGN.CENTER)
notes(s, "Close on the reframed headline: not a null, an unresolvable-but-favourable coincidence.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
