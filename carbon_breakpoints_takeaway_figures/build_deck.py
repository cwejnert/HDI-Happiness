#!/usr/bin/env python3
"""What Bends the Curve? — Paris at ten. v3: narrative-beat structure in the
'Measured Where It Varies' voice: beat kickers, sentence titles, prose not
bullets, act openers with stat callouts, and where-this-leaves-us interstitials."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

MEDIA = "/tmp/claude-0/-home-user-HDI-Happiness/1513cf0d-fca1-52a3-b470-357bab9ee4f8/scratchpad/deck_media/ppt/media"
FIGS = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/figures"
OUT = "/tmp/claude-0/-home-user-HDI-Happiness/1513cf0d-fca1-52a3-b470-357bab9ee4f8/scratchpad/what_bends_the_curve_final.pptx"

INK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x1B, 0x78, 0x37)
BLUE = RGBColor(0x3A, 0x66, 0xA5)
RED = RGBColor(0xB3, 0x3A, 0x3A)
AMBER = RGBColor(0xC7, 0x7F, 0x00)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         space_after=6, align=PP_ALIGN.LEFT, font="Calibri"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return p


def runs(tf, parts, size=14, first=False, space_after=6, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    for t, c, b in parts:
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, "Calibri"
        r.font.italic = italic
        r.font.color.rgb = c
    return p


def kicker_title(s, kicker, title, kcolor=GREEN):
    tf = box(s, 0.55, 0.3, 12.25, 1.0)
    para(tf, kicker, size=12, color=kcolor, bold=True, first=True, space_after=2)
    para(tf, title, size=24, color=INK, bold=True, space_after=0)


def fit_image(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih)
    w, h = iw * sc, ih * sc
    return s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2),
                                Inches(w), Inches(h))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def figure_slide(kicker, title, img, narrative=None, guide=None, note="", kcolor=GREEN):
    """Beat kicker, sentence title, one prose paragraph, an italic reading line, figure."""
    s = slide()
    kicker_title(s, kicker, title, kcolor=kcolor)
    top = 1.30
    if narrative:
        tf = box(s, 0.55, top, 12.25, 0.80)
        para(tf, narrative, size=13, color=GRAY, first=True, space_after=0)
        top += 0.84
    if guide:
        tf = box(s, 0.55, top, 12.25, 0.5)
        runs(tf, [("How to read it: ", BLUE, True), (guide, LGRAY, False)],
             size=11.5, first=True, space_after=0, italic=True)
        top += 0.52
    fit_image(s, img, 0.35, top + 0.04, 12.63, 7.2 - top)
    notes(s, note)
    return s


def act_open(roman, act_label, title, lede, stats, note=""):
    s = slide()
    tf = box(s, 0.9, 0.85, 11.5, 2.6)
    para(tf, roman, size=44, bold=True, color=GREEN, first=True, space_after=0)
    para(tf, act_label, size=12, bold=True, color=LGRAY, space_after=6)
    para(tf, title, size=28, bold=True, space_after=8)
    para(tf, lede, size=14, color=GRAY, space_after=0)
    x = 0.9
    for big, small in stats:
        tf = box(s, x, 4.6, 3.8, 2.0)
        para(tf, big, size=30, bold=True, color=GREEN, first=True, space_after=2)
        para(tf, small, size=12.5, color=GRAY, space_after=0)
        x += 3.95
    notes(s, note)
    return s


def interstitial(act_label, text, note=""):
    s = slide()
    tf = box(s, 0.9, 2.5, 11.5, 2.6)
    para(tf, act_label, size=13, bold=True, color=GREEN, first=True, space_after=10)
    para(tf, text, size=19, color=INK, space_after=0)
    notes(s, note)
    return s


M = lambda n: os.path.join(MEDIA, n)
F = lambda n: os.path.join(FIGS, n)

# ------------------------------------------------------------ 1 · title
s = slide()
tf = box(s, 0.9, 1.6, 11.5, 3.7)
para(tf, "What Bends the Curve?", size=46, bold=True, first=True, space_after=8)
para(tf, "An argument in four acts, on why the global carbon-intensity record cannot yet hand Paris the credit it may deserve — and where the policy signal actually lives.",
     size=19, color=GRAY, space_after=14)
runs(tf, [("Structural breaks · Kaya pathways · five decades of national decarbonization", GREEN, True)], size=14)
tf = box(s, 0.9, 5.95, 11.5, 0.9)
para(tf, "204 jurisdictions · 97 eligible national breaks · 107 transition episodes · 7 mechanism families · 9 constructive decarbonizations",
     size=13, color=INK, bold=True, first=True)
notes(s, "Frame up front: a pro-Paris paper that takes measurement seriously. The question is not whether Paris matters but what the observed record can attribute — and where the policy signal actually lives.")

# ------------------------------------------------------------ 2 · the setup
s = slide()
kicker_title(s, "THE SETUP", "Emissions growth has nearly stopped. Paris may not be why.")
tf = box(s, 0.55, 1.55, 7.35, 5.5)
para(tf, "The good news is real. Worst-case emissions scenarios have moved off the table, and global CO2 growth has slowed from +2.1% per year over the Rio-to-Paris quarter century to +0.4% per year since 2015. Something is going right.", size=14, color=GRAY, first=True, space_after=12)
para(tf, "But there is an obvious way to over-read that success: attribute the improvement to the agreement whose anniversary it coincides with. The world's carbon intensity has been falling for fifty years, for reasons that predate every climate treaty. So the question this paper asks is narrow and awkward: can the observed record distinguish a Paris effect from the continuation of trends already in motion?", size=14, color=GRAY, space_after=12)
para(tf, "Rather than assuming treaty years are turning points, we let the data say when trends actually changed — for the world, and for every country with enough annual data — and then ask the historical record why.", size=14, color=INK, space_after=0)
tf = box(s, 8.25, 1.55, 4.55, 5.5)
para(tf, "Why intensity, not emissions", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "Emissions arithmetic has three moving parts: economic growth, the fuel mix, and energy per dollar of output. Growth pushes emissions up; the two intensity terms are the only levers that push them down.", size=13, color=GRAY, space_after=10)
para(tf, "That makes intensity the leading indicator. A policy signal must appear in the intensity slopes before it can ever appear in emissions — and unless intensity improves faster than the economy grows, emissions cannot stabilize at all.", size=13, color=INK, space_after=0)
notes(s, "The sympathetic frame first, then the identity. The identity justifies the paper's design: the mathematically required location of any policy signal is the intensity slopes.")

# ------------------------------------------------------------ 3 · the argument
s = slide()
kicker_title(s, "THE ARGUMENT", "Four acts")
acts = [("I", "Did Paris bend the global curve?", "The global bends predate the climate regime; a new acceleration is real but began before 2015. The record shows the acceleration — not yet its author."),
        ("II", "What bends national curves?", "97 eligible breaks scatter across five decades and track economic history. Historical mechanisms, coded against documents, do the explaining."),
        ("III", "Did the bends count?", "A favorable statistic is not a transition. Four tests — direction, persistence, absolute CO2, mechanism — leave nine unambiguous successes."),
        ("IV", "Can we trust the story?", "Event timing beats chance; external raters agree; accounting choices bound the claims; and 22 unresolved cases stay unresolved.")]
x = 0.55
for roman, t, d in acts:
    tf = box(s, x, 1.7, 2.95, 4.8)
    para(tf, roman, size=38, bold=True, color=GREEN, first=True, space_after=6)
    para(tf, t, size=15.5, bold=True, space_after=8)
    para(tf, d, size=12.5, color=GRAY, space_after=0)
    x += 3.12
notes(s, "Act I is the reframed one: not 'treaties failed' but 'a real acceleration whose timing cannot yet separate Paris from the renewables cost decline'.")

# ------------------------------------------------------------ 4 · the instrument
s = slide()
kicker_title(s, "THE INSTRUMENT", "Two straight lines and a kink", kcolor=BLUE)
tf = box(s, 0.55, 1.55, 6.1, 5.5)
para(tf, "Every series is plotted in logs, so a straight line means a steady percentage change per year. We then fit two straight lines that meet at a kink, trying every admissible year and keeping the one that fits best. That year — the breakpoint — is the data's own estimate of when the trend changed, chosen with no reference to any treaty.", size=13.5, color=GRAY, first=True, space_after=12)
para(tf, "The date is then stress-tested: re-estimated on many statistically reshuffled versions of the series. A break whose year barely moves is credible; one that jumps around is reported as unstable. And throughout, a breakpoint answers only when a trend changed — never, by itself, why. Historical evidence answers that separately.", size=13.5, color=GRAY, space_after=12)
para(tf, "(Formally: y(t) = a + b·t + d·max(t−k, 0). Uncertainty via moving-block bootstrap; autocorrelation-robust standard errors for fixed-date tests; BIC for model evidence.)", size=10.5, color=LGRAY, italic=True, space_after=0)
tf = box(s, 7.0, 1.55, 5.8, 5.5)
para(tf, "Two doors for any improvement", size=15, bold=True, first=True, space_after=6)
para(tf, "Carbon intensity — CO2 per dollar of output — can only improve through one of two doors: cleaner energy (less CO2 per unit of energy: fuel switching, nuclear, renewables, the power system), or less energy per dollar (efficiency, and changes in what the economy makes).", size=13.5, color=GRAY, space_after=10)
para(tf, "We estimate the bend in the overall series and behind both doors, for the world and for every country, with one identical procedure. Which door moved constrains what happened: a fuel-mix bend points to energy policy or markets; an efficiency bend, as often as not, points to economic change.", size=13.5, color=INK, space_after=0)
notes(s, "Keep it conversational: two lines and a kink, tried at every year; then shake the data and see if the date holds still. The two-doors framing carries the mechanism logic of Acts II-III.")

# ------------------------------------------------------------ 5 · rules of evidence
s = slide()
kicker_title(s, "THE RULES OF EVIDENCE", "Most candidate breaks fail; treaty years face a harder test still", kcolor=BLUE)
tf = box(s, 0.55, 1.55, 5.9, 5.5)
para(tf, "Which breaks count", size=15, bold=True, first=True, space_after=6)
para(tf, "A country enters the analysis only when its series is long enough to be serious — forty-plus years, with at least ten on each side of the break — and only when the two-line model beats a single straight line decisively, the estimated year stays put under reshuffling, and the change in slope is large enough to matter rather than a statistical whisker.", size=13, color=GRAY, space_after=10)
para(tf, "Those rules take 204 jurisdictions down to 158 with estimable series, and to 97 countries with a break that passes every test. Breaks close together in time are read as one national transition — 107 episodes in all, so Bangladesh's cluster of 1990s breaks is one story, not three.", size=13, color=INK, space_after=0)
tf = box(s, 6.85, 1.55, 5.95, 5.5)
para(tf, "When a treaty year gets credit", size=15, bold=True, first=True, space_after=6)
para(tf, "The catch with treaty-year tests is that a gently curving trend makes almost any year look statistically significant. So a treaty year earns credit only if the slope genuinely changes there by a robust test, and the treaty year beats nearly all nearby placebo years, and the best-fitting date lands within three years of it, and the data's own chosen break agrees, and a documented domestic policy story fits.", size=13, color=GRAY, space_after=10)
para(tf, "This ladder is deliberately hard to climb: many series pass the first step, few pass all five. In Act I we hold Paris itself to the same standard.", size=13, color=INK, bold=True, space_after=0)
notes(s, "The placebo idea in one line: significance at the treaty year means little if 1990 and 1994 are just as significant. The ladder is the paper's discipline — and Act I turns it on Paris.")

# ------------------------------------------------------------ 6 · ACT I opener
act_open("I", "ACT I", "Did Paris bend the global curve?",
         "The global record shows two things at once: bends that long predate the climate regime, and a genuine new acceleration whose start the data cannot pin to 2015.",
         [("1973", "the year the world's carbon-intensity curve first bent"),
          ("+2.1% → +0.4%", "global CO2 growth per year, Rio-to-Paris era versus after Paris"),
          ("2013", "the best-fitting onset of the new acceleration — before the treaty")],
         note="Three numbers carry Act I: the old bend, the slowdown, and the awkward onset date.")

# ------------------------------------------------------------ 7 · the long trend
figure_slide("ACT I · THE LONG TREND",
             "The world's curve bent before the world began negotiating",
             F("takeaway_fig01_global_bend_predates_treaties.png"),
             narrative="The best-fitting break in global CO2 per dollar of output is 1973 — the oil-shock era, two decades before Rio and four before Paris. The demand side bent the same year. Only the fuel mix has an estimate near a treaty year, and its uncertainty interval is honest to the point of self-effacement: 1978 to 2011.",
             guide="Gray dots are the world's CO2 per dollar each year (1965 = 100, log scale); the red line is the best-fitting pair of straight lines, meeting at the estimated break. The lower panel gives each component's break year with its uncertainty whisker.",
             note="Robust to consumption-based emissions, excluding former Soviet economies, PPP GDP. Multiple-break models show several episodes, not one intervention.")

# ------------------------------------------------------------ 8 · anniversary test
figure_slide("ACT I · THE ANNIVERSARY TEST",
             "Ten years on, the world is ahead of its old trend",
             F("takeaway_fig02a_paris_and_the_long_trend.png"),
             narrative="Extend the Rio-to-Paris trend past 2015 and the years since fall below it — a real acceleration, concentrated where climate policy operates: the fuel mix. But onsets from 2011 to 2016 fit about equally well (2013 best), and the demand side shows no acceleration. The world is doing what an effective Paris requires; the record cannot yet say Paris is why.",
             guide="The solid line is the 1990–2015 trend; the dashed line extends it with a band for where years should fall if nothing changed. Red dots are the years since Paris — below the band means faster-than-trend improvement.",
             note="C/E −0.7 pp/yr (NW t=4.3); C/GDP −0.9 (t=2.9); E/GDP −0.3 (t=1.5, ns). Onset plateau 2011–2016, best 2013 — indistinguishable from the renewables cost collapse.")

# ------------------------------------------------------------ 9 · the arithmetic
figure_slide("ACT I · THE ARITHMETIC",
             "Emissions bend only when intensity outruns growth — it doesn't, yet",
             F("takeaway_fig02b_stabilization_arithmetic.png"),
             narrative="By construction, emissions growth equals GDP growth plus the two intensity terms — so the dotted line, where intensity decline equals GDP growth, is the threshold for emissions to stop rising. The slowdown since Paris is real, but roughly two-thirds of it is slower economic growth; one-third is the cleaner fuel mix. The gap to the dotted line is the remaining task.",
             guide="Each bar is an average annual growth rate; the red bar equals the gray bar plus the two colored ones.",
             note="Emissions +2.1 to +0.4%/yr after 2015. Decomposition: ~1.1pp slower GDP growth (incl. COVID), ~0.5pp fuel mix, ~0.1pp efficiency.")

# ------------------------------------------------------------ 10 · placebo discipline
figure_slide("ACT I · THE PLACEBO DISCIPLINE",
             "Any smooth curve makes many years look significant",
             F("takeaway_fig02_placebo_test.png"),
             narrative="Place the break at every candidate year and Rio fits well — alongside a broad plateau of neighbors that fit as well or better. That is why treaty-year significance, common in half to two-thirds of country series, collapses to a small minority once each treaty year must beat its neighbors. The same discipline, turned on Paris's decade, dates the new acceleration to roughly 2013.",
             guide="Top: model improvement from placing the global fuel-mix break in each year. Bottom: the share of country series significant at each treaty date (left point) versus the share where that date also beats nearby placebo years (right point).",
             note="Rio survivors: 10 C/GDP, 21 C/E, 34 E/GDP series; Paris: 14, 5, 6. The Rio-unique E/GDP cases dissolve on inspection into early-90s upheaval — disruption, adjustment, liberalization — not treaty policy.")

# ------------------------------------------------------------ 11 · where act I leaves us
s = slide()
kicker_title(s, "WHERE ACT I LEAVES US", "The record shows an acceleration — not yet its author")
tf = box(s, 0.55, 1.6, 6.05, 5.4)
para(tf, "Established", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "A real post-2013 supply-side acceleration — the fuel mix improving at −0.5% per year after a flat quarter-century — and emissions growth down from +2.1% to +0.4% per year. This is the direction an effective Paris requires, in the component where policy operates most directly. The demand side, meanwhile, continues its fifty-year trend, no faster.", size=13.5, color=GRAY, space_after=0)
tf = box(s, 7.0, 1.6, 5.8, 5.4)
para(tf, "Open — and why", size=15, bold=True, color=RED, first=True, space_after=6)
para(tf, "The onset is timed to ~2013, not 2015, and eight post-Paris years give the test almost no power: Paris would have had to double a fifty-year trend to prove itself already. The world series is also a China-weighted average — the WTO-era coal boom flattened the baseline that makes the recent years look fast. And Paris was built to work through national pledges and national transformation, so its signature should appear country by country.", size=13.5, color=GRAY, space_after=8)
para(tf, "So we look there.", size=15, color=INK, bold=True, space_after=0)
notes(s, "The hinge of the talk: the non-attribution is not a failure of Paris — it is the reason the analysis must descend to the national level, where attribution has power.")

# ------------------------------------------------------------ 12 · ACT II opener
act_open("II", "ACT II", "What bends national curves?",
         "If the global aggregate cannot separate treaty from trend, the country record can at least say what bends carbon-intensity curves when they do bend — and it is rarely a treaty date.",
         [("97", "national breaks passing every statistical screen, across five decades"),
          ("45% vs 17%", "documented events within two years of a break, observed versus chance"),
          ("9 of 97", "breaks best explained by climate policy and low-carbon technology")],
         note="Act II is deliberately light-touch in this telling: the scatter, the census, then straight to whether the bends counted.")

# ------------------------------------------------------------ 13 · the scatter
figure_slide("ACT II · THE SCATTER",
             "Ninety-seven breaks, five decades, no treaty era",
             F("takeaway_fig03_when_and_why_curves_bend.png"),
             narrative="National breaks pile up around the oil shocks, the post-Cold-War upheaval, the Asian and global financial crises, and the renewables era — around economic history, not around the two treaty lines. And the clustering is not storytelling after the fact: documented events sit within two years of the estimated breaks nearly three times as often as randomly drawn dates would allow.",
             guide="Top: one dot per country at its break year; green triangles are the nine policy-enabled cases; shaded bands mark major economic episodes. Bottom: observed event-to-break proximity (red) against what random timing produces (blue).",
             note="45% vs 17% within 2 years; 80% vs 37% within 5; 70% vs 33% inside the documented episode.")

# ------------------------------------------------------------ 14 · the census
figure_slide("ACT II · THE CENSUS",
             "Mostly markets, restructuring, and disruption — policy in nine",
             F("takeaway_fig04_mechanism_census.png"),
             narrative="Each break is assigned the historical process best supported by documentary evidence: fuel and energy markets lead with 22 countries, economic restructuring takes 18, political disruption 11, macroeconomic shocks 9, development and energy access 6. Climate policy and low-carbon technology account for nine — real, but rare. Twenty-two breaks stay honestly unresolved rather than being forced into a category.",
             guide="One square per country, grouped by best-supported mechanism; green is the policy group, light gray the unresolved.",
             note="Author-coded against authoritative sources with an explicit stopping rule — structured interpretations, not causal estimates.")

# ------------------------------------------------------------ 15 · ACT III opener
act_open("III", "ACT III", "Did the bends count?",
         "Every one of the 97 breaks looks favorable as a statistic. Four tests — did the trend improve, did it last, did absolute CO2 fall, was the mechanism constructive — sort statistics from transitions.",
         [("55", "countries whose intensity improved while absolute CO2 kept rising"),
          ("20", "whose 'improvement' is a recession or collapse in disguise"),
          ("9", "constructive, persistent decarbonizations — the gold standard")],
         note="The taxonomy slide follows as the reading key; then the flow, the quadrant, the nine, and the near-misses.")

# ------------------------------------------------------------ 16 · the sorting key
s = slide()
kicker_title(s, "ACT III · THE SORTING KEY", "Seven mechanisms in, six verdicts out", kcolor=AMBER)
tf = box(s, 0.55, 1.5, 6.05, 5.7)
para(tf, "MECHANISMS — what accompanied the bend", size=14, bold=True, first=True, space_after=6)
for name, d in [("Policy-enabled low carbon", "climate/energy policy made physical: infrastructure, fuel mix, market design"),
                ("Fuel & energy markets", "market-driven fuel switching or supply transformation"),
                ("Economic restructuring", "industry to services, liberalization, post-socialist transition"),
                ("Macroeconomic shock", "recessions and financial crises that cut energy use"),
                ("Political disruption", "wars, state collapse, sanctions, regime change"),
                ("Development & energy access", "electrification and modernization"),
                ("No attributed mechanism", "a real break the record cannot yet explain (22 of 97)")]:
    runs(tf, [(name + " — ", INK, True), (d + ".", GRAY, False)], size=12, space_after=4)
tf = box(s, 7.0, 1.5, 5.8, 5.7)
para(tf, "VERDICTS — four tests of the outcome", size=14, bold=True, first=True, space_after=4)
para(tf, "Did the trend improve? Did it last? Did absolute CO2 actually fall? Was the underlying process constructive?", size=12.5, italic=True, color=GRAY, space_after=8)
for name, n, d in [("Constructive persistent decarbonization", 9, "passes all four — the gold standard"),
                   ("Persistent fuel substitution", 13, "durable supply-side switch; narrower transformation"),
                   ("Intensity gain; CO2 still rising", 12, "relative gain overwhelmed by growth"),
                   ("Contractionary / disruptive", 20, "the 'improvement' is a crisis in disguise"),
                   ("Incomplete / nonpersistent", 21, "the bend reversed or did not hold"),
                   ("Outcome not classifiable", 22, "too little post-break evidence")]:
    runs(tf, [(f"{name} ", INK, True), (f"({n}) — ", GREEN if n == 9 else LGRAY, True), (d + ".", GRAY, False)], size=12, space_after=4)
notes(s, "The reading key for everything that follows. Note the two distinct 22s: mechanism-unresolved and outcome-unclassifiable are different sets.")

# ------------------------------------------------------------ 17 · the flow
figure_slide("ACT III · THE FLOW",
             "Follow any mechanism to where its transitions actually end up",
             M("image4.png"),
             narrative="Shocks and disruption flow into contraction. Restructuring most often ends in intensity gains that growth overwhelms. The policy-enabled ribbon is thin — and it is the one that most reliably reaches the constructive verdict on the right.",
             guide="Each ribbon carries countries from the mechanism behind their break (left) to the verdict on the outcome (right); width is the number of countries.",
             note="Only nine countries end in constructive persistent decarbonization.", kcolor=AMBER)

# ------------------------------------------------------------ 18 · the quadrant
figure_slide("ACT III · THE QUADRANT",
             "Lower intensity is not the same as falling emissions",
             F("takeaway_fig05_bend_vs_success.png"),
             narrative="Fifty-five countries improved their intensity trend while absolute emissions kept rising; nineteen managed both. The nine policy-enabled cases, marked in green, mostly sit where success lives — but Bhutan and Paraguay improve intensity with emissions still rising, which is exactly why 'policy-enabled' and 'constructive' are overlapping but different sets.",
             guide="Each dot is a country: further right, a bigger post-break improvement in the intensity trend; below the dashed line, absolute CO2 actually falling. Success lives only in the lower right.",
             note="Raw counts 55/19; the stricter classification (adding persistence and mechanism) narrows these to 12 and 9.", kcolor=AMBER)

# ------------------------------------------------------------ 19 · the nine
figure_slide("ACT III · THE NINE",
             "Nine countries pass all four tests",
             M("image7.png"),
             narrative="Sweden 1986 · Finland 1993 · Hungary 1996 · Switzerland 1999 · Nauru 1999 · Denmark 2002 · Portugal 2004 · United Kingdom 2010 · Ireland 2012. Six of nine are policy-enabled; Hungary, Nauru, and Ireland arrived via restructuring — and Nauru's is a phosphate collapse, so it carries an asterisk.",
             guide="Each panel sets both series to 100 at the country's break year: blue is carbon intensity, red is absolute CO2. The defining test is that both fall together after the break (shaded region).",
             note="This is the paper's constructive core: policy made physical in fuel mix, infrastructure, and economic structure.", kcolor=AMBER)

# ------------------------------------------------------------ 20 · the near misses
figure_slide("ACT III · THE NEAR MISSES",
             "Nineteen countries cut emissions; ten don't earn the verdict",
             M("image8.png"),
             narrative="Falling emissions alone are not enough — the route matters. Five of the nineteen are durable fuel substitutions: the Netherlands did with gas in 1975 what the US did with shale forty years later. Somalia 1985 is the caution: collapse mimics decarbonization in every raw statistic; only the persistence screen catches it.",
             guide="Each row is a country where absolute CO2 fell after a favorable bend, with its break-year uncertainty whisker, grouped by verdict; labels name the best-supported mechanism.",
             note="Even the genuine emission-cutters span 1975–2012 — no treaty era here either.", kcolor=AMBER)

# ------------------------------------------------------------ 21 · the map
s = slide()
kicker_title(s, "ACT III · THE GEOGRAPHY", "Mechanisms are regional stories; success is concentrated", kcolor=AMBER)
tf = box(s, 0.55, 1.3, 12.25, 0.72)
para(tf, "Fuel-market transformations cluster in the Americas and the Middle East; restructuring runs across Asia-Pacific; the policy-enabled cases concentrate in northern Europe. On the verdict map, most of the world is contraction, incompleteness, or intensity gains that growth overwhelms — the dark-green constructive nine occupy one corner of one continent.",
     size=13, color=GRAY, first=True, space_after=0)
fit_image(s, M("image9.png"), 0.30, 2.1, 6.35, 5.0)
fit_image(s, M("image10.png"), 6.75, 2.1, 6.35, 5.0)
notes(s, "Left map: mechanism per country. Right map: outcome verdict; dark green outlines the nine.")

# ------------------------------------------------------------ 22 · ACT IV opener
act_open("IV", "ACT IV", "Can we trust the story?",
         "Three checks that could have broken the argument, and one set of limits stated plainly.",
         [("80% vs 37%", "documented events within five years of breaks, observed versus chance"),
          ("64th vs 34th", "median external climate-rating percentile, policy-enabled cases versus the rest"),
          ("6 of 17", "countries whose break date survives switching to consumption-based accounting")],
         note="Act IV: validation that carries temporal information; corroboration on thin coverage; accounting humility.")

# ------------------------------------------------------------ 23 · corroboration
figure_slide("ACT IV · THE CORROBORATION",
             "Independent raters recognize the same countries",
             F("takeaway_fig06_external_corroboration.png"),
             narrative="Climate Action Tracker and the Climate Change Performance Index know nothing of our breakpoints, yet the policy-enabled cases land high on both — median percentiles of 64 versus 34, and 91 versus 44. Coverage is the honest caveat: only three of the nine appear in CAT and two in CCPI, so every covered country is shown rather than summarized away.",
             guide="Each dot is a covered country placed at its rating percentile; green triangles are the policy-enabled cases; vertical bars mark group medians.",
             note="Corroborative, not definitive.", kcolor=RED)

# ------------------------------------------------------------ 24 · limits
s = slide()
kicker_title(s, "ACT IV · THE LIMITS", "What the design supports, and what it refuses to claim", kcolor=RED)
col = [("Consumption-based accounting",
        "Assign traded emissions to consumers instead of producers and only 6 of 17 comparable countries keep a break date within five years; the median shift is twelve years. Some 'domestic' transitions partly reflect offshoring. This is a main result, not a footnote."),
       ("Carbon pricing & stringency",
        "Carbon pricing is associated with a cleaner fuel mix in panel models with country and year controls; the other components are imprecise, and stringency data cover only 15 of 28 focal countries. Associational, not causal — and said so."),
       ("Event-timing validation",
        "The historical classifications carry real temporal information: documented events fall within two years of the estimated break in 45% of episodes against ~17% by chance, and inside the documented episode in 70% against 33%."),
       ("What breakpoints can't say",
        "Break timing identifies change, not cause. Attribution and outcome classification each leave 22 cases open — restraint that makes the resolved three-quarters credible. And a treaty that prevents backsliding produces no break at all: absence of a break is not absence of an effect.")]
xs, ys = [0.55, 7.0, 0.55, 7.0], [1.55, 1.55, 4.3, 4.3]
for (t, d), x, y in zip(col, xs, ys):
    tf = box(s, x, y, 5.9, 2.65)
    para(tf, t, size=14.5, bold=True, first=True, space_after=4)
    para(tf, d, size=12, color=GRAY, space_after=0)
notes(s, "The credibility slide. If pressed on causality, lean on the design lock's explicit list of supported and unsupported claims.")

# ------------------------------------------------------------ 25 · the verdict
figure_slide("THE VERDICT",
             "Curves bend everywhere; decarbonization is rare",
             M("image12.png"),
             narrative="All 97 national breaks looked favorable as statistics. Checked for persistence, mechanism, and absolute emissions, they resolve into a spectrum from crisis-in-disguise to genuine transformation — and none of this heterogeneity aligns on a treaty date. The treaty question dissolves into 97 national stories, nine of which end in unambiguous decarbonization.",
             guide="One block per verdict, sized by countries; only the dark-green block on the right is unambiguous decarbonization.",
             note="The bridge to the close: heterogeneous national mechanisms, not one universal treaty break.")

# ------------------------------------------------------------ 26 · the ask
s = slide()
kicker_title(s, "THE ASK", "Judge Paris where its signal must appear — and check again at twenty")
tf = box(s, 0.55, 1.7, 11.9, 4.8)
para(tf, "If the argument holds, the right test of the climate regime is not whether emissions charts bend at signing ceremonies. It is whether national intensity slopes accelerate, country by country, in ways that placebo years cannot mimic and documented policy can explain — the test the nine already pass.", size=15.5, color=GRAY, first=True, space_after=12)
para(tf, "On that test, Paris at ten is neither vindicated nor indicted: the world is finally moving the way an effective agreement would require, for reasons the record cannot yet separate from the technology revolution the agreement helped consolidate. Paris at twenty will be decidable — if the acceleration spreads from the fuel mix to the demand side, and from nine countries to many.", size=15.5, color=GRAY, space_after=12)
para(tf, "The question we are putting to the field: should treaty assessment be built on intensity breakpoints and placebo discipline, rather than on emissions levels and anniversaries?", size=15.5, color=INK, bold=True, space_after=0)
notes(s, "The forward-looking beat: what would count as evidence at Paris+20, and the methods question posed to the field.")

# ------------------------------------------------------------ 27 · close
s = slide()
tf = box(s, 0.9, 1.5, 11.5, 1.9)
para(tf, "Global agreements can organize climate action —", size=27, bold=True, first=True, space_after=2)
para(tf, "durable national transformations bend the curve.", size=27, bold=True, color=GREEN, space_after=0)
tf = box(s, 0.9, 3.5, 11.5, 2.9)
for t in ["Global curves bent in the 1970s; since ~2013 they bend faster — a real supply-side acceleration whose timing cannot yet separate Paris from the renewables revolution it helped consolidate.",
          "Favorable bends often reflect recession, disruption, deindustrialization, or offshoring; intensity gains are routinely overwhelmed by growth.",
          "Nine countries pair persistent intensity declines with falling absolute CO2 — policy made physical in fuel mix, infrastructure, and economic structure."]:
    runs(tf, [("—  ", GREEN, True), (t, GRAY, False)], size=14.5, space_after=9)
tf = box(s, 0.9, 6.55, 11.5, 0.7)
para(tf, "Structural breaks say when · Kaya says through which pathway · historical validation says what plausibly happened · transition quality says whether it mattered",
     size=13, italic=True, color=LGRAY, first=True, align=PP_ALIGN.CENTER)
notes(s, "End on the closing line and the method signature.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
