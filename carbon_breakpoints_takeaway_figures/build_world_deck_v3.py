#!/usr/bin/env python3
"""Paris at Ten — v3. Rebuilt on a true world aggregate and a 44-specification grid."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

FIGS = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/outputs/final_master/figures_world_v3"
OUT = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/paris_at_ten_world_deck.pptx"

INK = RGBColor(0x21, 0x21, 0x21); GRAY = RGBColor(0x55, 0x55, 0x55); LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x1B, 0x78, 0x37); BLUE = RGBColor(0x3A, 0x66, 0xA5)
RED = RGBColor(0xB3, 0x3A, 0x3A); AMBER = RGBColor(0xC7, 0x7F, 0x00)

prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
slide = lambda: prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         space_after=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, "Calibri"
    r.font.color.rgb = color
    return p


def runs(tf, parts, size=14, first=False, space_after=6, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    for t, c, b in parts:
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, "Calibri"
        r.font.italic = italic; r.font.color.rgb = c
    return p


def kicker_title(s, kicker, title, kcolor=GREEN):
    tf = box(s, 0.55, 0.3, 12.25, 1.0)
    para(tf, kicker, size=12, color=kcolor, bold=True, first=True, space_after=2)
    para(tf, title, size=24, color=INK, bold=True, space_after=0)


def fit_image(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih); w, h = iw * sc, ih * sc
    return s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2), Inches(w), Inches(h))


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


def finding(kicker, title, takeaway, img, caveat=None, note="", kcolor=GREEN):
    s = slide(); kicker_title(s, kicker, title, kcolor=kcolor)
    tf = box(s, 0.55, 1.30, 12.25, 0.5)
    runs(tf, [("Takeaway:  ", kcolor, True), (takeaway, INK, True)], size=14, first=True, space_after=0)
    top = 1.86
    if caveat:
        tf = box(s, 0.55, top, 12.25, 0.42)
        runs(tf, [("But:  ", RED, True), (caveat, GRAY, False)], size=12, first=True, space_after=0)
        top += 0.46
    fit_image(s, img, 0.35, top, 12.63, 7.2 - top)
    notes(s, note); return s


F = lambda n: os.path.join(FIGS, n)

# 1 · title
s = slide()
tf = box(s, 0.9, 1.7, 11.5, 3.6)
para(tf, "Paris at Ten", size=46, bold=True, first=True, space_after=8)
para(tf, "The world's carbon intensity is falling faster than ever. We ask whether the record can credit the Paris Agreement — and find that it cannot, for a specific and instructive reason.",
     size=18, color=GRAY, space_after=14)
runs(tf, [("True world aggregate, 1965-2023 · 44 specifications · unknown-break-date tests", GREEN, True)], size=14)
notes(s, "World-level paper. The country analysis is the follow-on. Every number here is from the true world aggregate unless stated; the specification grid is the robustness spine.")

# 2 · the question
s = slide()
kicker_title(s, "THE QUESTION", "Something changed. Was it Paris?")
tf = box(s, 0.55, 1.6, 7.4, 5.4)
para(tf, "The good news is real. Global CO2 growth fell from +2.1% a year over the Rio-to-Paris quarter century to +0.7% a year since, and carbon intensity is now improving at -2.1% a year: the fastest sustained rate in the observed record.",
     size=15, color=GRAY, first=True, space_after=14)
para(tf, "But the obvious reading — that the agreement did it — is the reading the data is worst at supporting. Carbon intensity has been falling for fifty years, in episodes with their own causes. Clean-energy costs collapsed on almost the same schedule as the treaty. And emissions growth slows when the economy slows.",
     size=15, color=GRAY, space_after=14)
para(tf, "So we ask a narrow question: when did the trend actually turn, and is 2015 a special year in the record — or a year that happens to sit near one?",
     size=15, color=INK, bold=True, space_after=0)
tf = box(s, 8.35, 1.6, 4.45, 5.4)
para(tf, "Why intensity, not emissions?", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "Policy signals show up first in intensity. Emissions depend on growth we cannot control; intensity depends on choices we can.",
     size=13.5, color=GRAY, space_after=8)
para(tf, "The decomposition: C/GDP = (C/E) × (E/GDP)", size=13, color=INK, bold=True, space_after=8)
para(tf, "Fuel mix and efficiency are where policy operates. A real policy effect must show up in one before the attribution is credible.",
     size=13.5, color=GRAY, space_after=0)
notes(s, "Policy can affect fuel mix (technology, prices) and efficiency (demand, standards), but not growth directly. Intensity strips out growth noise.")

# 3 · what we did
s = slide()
kicker_title(s, "WHAT WE DID", "Four steps", kcolor=BLUE)
steps = [("1", "Start with real data",
          "The true world totals (OWID CO2 and energy, World Bank GDP), not sums of whatever countries reported. This avoids the problem where adding more countries changes the baseline."),
         ("2", "Find the turning point",
          "Fit two straight-line segments and try every possible break year. Keep the one that fits best. Use statistical testing to make sure we're not just seeing random noise."),
         ("3", "Check if Paris is special",
          "Any year can look significant if the trend is curving. We rank 2015 against every other year — Paris only gets credit if 2015 is genuinely more significant than 2014 or 2016."),
         ("4", "Test robustness",
          "Run the whole analysis 44 different ways: different starting years, endpoints, how much data we require, whether we fill gaps. Report only results that survive these tests.")]
x = 0.55
for n, t, d in steps:
    tf = box(s, x, 1.7, 3.0, 5.0)
    para(tf, n, size=32, bold=True, color=BLUE, first=True, space_after=6)
    para(tf, t, size=14.5, bold=True, space_after=8)
    para(tf, d, size=12, color=GRAY, space_after=0)
    x += 3.13
tf = box(s, 0.55, 6.7, 12.25, 0.6)
runs(tf, [("Key principle: a turning point shows ", INK, True),
          ("when", GREEN, True),
          (" a trend changed, never ", INK, True),
          ("why. ", GRAY, False)], size=13.5, first=True)
notes(s, "Step 4 is crucial: it catches overclaims. It also lets us understand what's truly robust versus what depends on how you build the data.")

# 3b · policy context (NEW SLIDE)
s = slide()
kicker_title(s, "THE STAKES", "Why intensity matters for climate")
tf = box(s, 0.55, 1.6, 12.25, 5.5)
para(tf, "Global carbon intensity is now declining at -2.1% per year — the fastest rate ever recorded.",
     size=15, bold=True, color=GREEN, first=True, space_after=12)
para(tf, "But there's a question underneath every climate target: is that fast enough?",
     size=15, color=GRAY, space_after=12)
para(tf, "Climate models tell us what rates are needed: about -2.5%/yr to stabilize emissions while the world grows, -5%/yr to start cutting absolute emissions, and -11%/yr to halve emissions by 2050.",
     size=14, color=GRAY, space_after=12)
para(tf, "We're currently above the stabilization target but far below what deep cuts require. The key question for policy: Is the 2012 acceleration real, will it persist, and can it accelerate further to meet climate goals?",
     size=14, color=INK, bold=True, space_after=0)
notes(s, "Context: current world intensity decline at -2.1%/yr. Stabilization needs -2.5%/yr. Halving emissions by 2050 needs -11%/yr. We have to understand whether what happened at 2012 is technology momentum or policy change, because that affects whether we can dial it up.")
finding("WHAT WE FIND · 1 OF 4", "The decisive turn in the record is the 1970s — not a treaty",
        "The oil shocks remain the largest structural break in the world's carbon intensity. The current acceleration began in 2012 — three years before Paris.",
        F("v3_fig1_long_record.png"),
        caveat="This tells us that long-term structural forces, not recent agreements, have been the largest drivers of intensity change.",
        note="Formally: unknown-break test shows dominant break at 1972 for overall intensity (p=0.002) and 1974 for efficiency (p<0.001). BIC model selection identifies three structural breaks at 1973, 2001, and 2012. The 2012 break is robust: it appears in 89% of specifications.")

# 5 · finding 2
finding("WHAT WE FIND · 2 OF 4", "Efficiency never changed pace. The fuel mix stalled for 25 years, then resumed.",
        "Energy efficiency improved at a steady 1% per year throughout the entire period — unchanged by Paris. The fuel mix stalled completely from 1990-2015, then resumed.",
        F("v3_fig2_baseline_problem.png"),
        caveat="This is the key policy result: efficiency (how much energy we use per dollar of GDP) didn't respond to the treaty. Only the fuel mix (power source composition) changed.",
        note="Fuel-mix rates: -0.52%/yr (oil shock recovery 1973-90), +0.02%/yr (stalled Rio-Paris), -0.70%/yr (resumed post-Paris). Efficiency unchanged: ~-1.0%/yr all three eras. The 25-year stall in fuels was the anomaly; recent resumption reverts to historical pace.")

# 6 · finding 3
finding("WHAT WE FIND · 3 OF 4", "The turn happens before Paris, no matter how we build the data",
        "Across 44 different ways of constructing the data, 89-93% place the turning point before 2015 — this is the timing question Paris attribution must answer.",
        F("v3_fig3_specification_invariance.png"),
        caveat="The turning point robustness is what makes the timing question meaningful. This is not about the magnitude of change, only when it began.",
        note="Across 44 specifications (different start years 1965-1990, endpoints 2020-2023, completeness thresholds, with/without interpolation): 89% for overall intensity, 93% for fuel mix turn before 2015. Typical break year: 2011-2013. This is the single most robust result.",
        kcolor=AMBER)

# 7 · finding 4
finding("WHAT WE FIND · 4 OF 4", "When we test everything, only timing is robust. Magnitude is not.",
        "The 2015 break survives testing, but whether the post-2015 change is statistically significant varies: 78% for fuel mix, but only 41% for efficiency. This is where data construction choices matter.",
        F("v3_fig4_what_is_robust.png"),
        caveat="Paris ranks 80th percentile for overall intensity but only 11th percentile for the fuel mix — it does not cleanly beat its neighbouring years.",
        note="Specification grid shows: timing robust (89-93% before 2015), magnitude fragile (78% find significant fuel-mix change, 41% for efficiency). This tells us the timing question is robust but the attribution question depends heavily on construction choices and which component you study.",
        kcolor=AMBER)

# 8 · verdict
s = slide()
kicker_title(s, "WHAT IT ADDS UP TO", "A real acceleration that began before the treaty")
tf = box(s, 0.55, 1.6, 6.05, 5.4)
para(tf, "What we can say", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "The world is decarbonizing faster than at any point in the record: carbon intensity has fallen at 2.1% a year since 2012, against 0.4% before. That improvement is real, large, and concentrated in the fuel mix — where technology and policy operate. Current pace is faster than the long-term trend.",
     size=13.5, color=GRAY, space_after=0)
tf = box(s, 7.0, 1.6, 5.8, 5.4)
para(tf, "What we cannot claim", size=15, bold=True, color=RED, first=True, space_after=6)
para(tf, "The turn dates to 2012 — three years before Paris. Energy efficiency never responded to the treaty at all. The fuel-mix improvement matches clean-energy cost collapse timing, not the Paris Agreement. Attribution to Paris is not supported by the timing.",
     size=13.5, color=GRAY, space_after=10)
para(tf, "The acceleration is real and necessary. It started before Paris. The evidence points to technology and cost, not yet to policy.",
     size=13.5, color=INK, bold=True, space_after=0)
notes(s, "Honest verdict: world is decarbonizing faster, but before the treaty and without efficiency gains. The acceleration looks like what technology-driven change produces.")

# 9 · what would settle it
s = slide()
kicker_title(s, "WHAT WOULD SETTLE IT", "Three questions for attribution")
tf = box(s, 0.55, 1.65, 11.9, 5.0)
para(tf, "Does it last? A 2012 turn from pure technology shock fades as cheap solar saturates. One from policy ratcheting persists. Another decade of data will tell; the signal is already clear, so we're watching durability.",
     size=15, color=GRAY, first=True, space_after=12)
para(tf, "Does efficiency wake up? Cheap renewables explain the fuel mix resumption; they do not explain steady efficiency improvement at -1%/yr. If efficiency accelerates post-Paris, Paris gains credibility. If it stays flat, technology-only remains the simpler explanation.",
     size=15, color=GRAY, space_after=12)
para(tf, "What do countries show? Paris works through national pledges and targets. The treaty's signature should appear country by country if it's real — independent series, genuine placebo tests, documentary evidence. The country-level work will answer this.",
     size=15, color=GRAY, space_after=14)
para(tf, "Core question for the field: should we evaluate treaties by intensity turning points and robustness testing, or by emissions trajectories and post-hoc explanations?",
     size=15, color=INK, bold=True, space_after=0)
notes(s, "These three questions distinguish technology-driven change from policy-driven change. The country paper will test the most important one.")

# 10 · limits
s = slide()
kicker_title(s, "LIMITS", "What we are not claiming", kcolor=RED)
col = [("This is not causal",
        "A turning point identifies when a trend changed, never why. Even a perfect fit at 2015 would be a coincidence of timing. And an agreement that prevented backsliding would leave no turning point at all — absence of one is not absence of an effect."),
       ("Two GDP concepts",
        "Overall intensity and efficiency use World Bank world GDP in constant 2015 dollars; the country panels use purchasing-power GDP. The fuel-mix result — the one most relevant to policy — uses a single source and no GDP at all."),
       ("Eight post-Paris years",
        "The series ends in 2023 for CO2 and energy and 2022 for GDP. Country-level work stops in 2021 because GDP reporting, not emissions reporting, runs out. Adding partial years would overstate the recent trend."),
       ("Production-based accounting",
        "Emissions are counted where they are produced. At the world level that is the right frame, since global production equals global consumption — but it matters for the national work that follows.")]
xs, ys = [0.55, 7.0, 0.55, 7.0], [1.6, 1.6, 4.35, 4.35]
for (t, d), x, y in zip(col, xs, ys):
    tf = box(s, x, y, 5.9, 2.6)
    para(tf, t, size=14.5, bold=True, first=True, space_after=4)
    para(tf, d, size=12, color=GRAY, space_after=0)
notes(s, "Volunteer the GDP-concept issue; a referee will find it.")

# 11 · close
s = slide()
tf = box(s, 0.9, 1.6, 11.5, 2.0)
para(tf, "The world is decarbonizing. It started in 2012, not 2015.", size=28, bold=True, first=True, space_after=2)
para(tf, "That's the robust finding. What comes next is uncertain.", size=24, bold=True, color=AMBER, space_after=0)
tf = box(s, 0.9, 3.7, 11.5, 2.9)
for t in ["The 1970s oil shocks remain the largest turn in the record — structural forces still matter most.",
          "The current acceleration started in 2012, before Paris, and in 89% of specifications.",
          "Efficiency never changed pace at all — policy has not yet moved the demand side of the equation.",
          "The fuel mix resumed its historical decline — matching the timing of renewable cost collapse."]:
    runs(tf, [("—  ", GREEN, True), (t, GRAY, False)], size=15, space_after=10)
tf = box(s, 0.9, 6.6, 11.5, 0.7)
para(tf, "What happened (timing) is clear · How (fuel mix) is clearer than demand · Why (technology or policy?) needs country-level tests · Whether it's enough is another question",
     size=12.5, italic=True, color=LGRAY, first=True, align=PP_ALIGN.CENTER)
notes(s, "The message: real acceleration, pre-treaty timing, technology-driven so far, country-level work coming. What it means for Paris depends on what the countries show.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
