#!/usr/bin/env python3
"""Paris at Ten — world-level deck, v2. Simpler: what we did, then what we find,
with one bold takeaway on every figure slide."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

FIGS = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/outputs/final_master/figures_world_v2"
OUT = "/home/user/HDI-Happiness/carbon_breakpoints_takeaway_figures/paris_at_ten_world_deck.pptx"

INK = RGBColor(0x21, 0x21, 0x21); GRAY = RGBColor(0x55, 0x55, 0x55); LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x1B, 0x78, 0x37); BLUE = RGBColor(0x3A, 0x66, 0xA5)
RED = RGBColor(0xB3, 0x3A, 0x3A); AMBER = RGBColor(0xC7, 0x7F, 0x00)

prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
slide = lambda: prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=14, color=INK, bold=False, italic=False, first=False,
         space_after=6, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.italic, r.font.name = Pt(size), bold, italic, "Calibri"
    r.font.color.rgb = color
    return p


def runs(tf, parts, size=14, first=False, space_after=6, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    for t, c, b in parts:
        r = p.add_run(); r.text = t
        r.font.size, r.font.bold, r.font.name = Pt(size), b, "Calibri"
        r.font.italic = italic; r.font.color.rgb = c
    return p


def kicker_title(s, kicker, title, kcolor=GREEN):
    tf = box(s, 0.55, 0.3, 12.25, 1.0)
    para(tf, kicker, size=12, color=kcolor, bold=True, first=True, space_after=2)
    para(tf, title, size=24, color=INK, bold=True, space_after=0)


def fit_image(s, path, x, y, maxw, maxh):
    iw, ih = Image.open(path).size
    sc = min(maxw / iw, maxh / ih); w, h = iw * sc, ih * sc
    return s.shapes.add_picture(path, Inches(x + (maxw - w) / 2), Inches(y + (maxh - h) / 2), Inches(w), Inches(h))


def notes(s, t):
    s.notes_slide.notes_text_frame.text = t


def finding(kicker, title, takeaway, img, caveat=None, note="", kcolor=GREEN):
    """Kicker, plain-language title, one bold takeaway, the figure, optional caveat line."""
    s = slide(); kicker_title(s, kicker, title, kcolor=kcolor)
    tf = box(s, 0.55, 1.30, 12.25, 0.5)
    runs(tf, [("Takeaway:  ", kcolor, True), (takeaway, INK, True)], size=14, first=True, space_after=0)
    top = 1.86
    if caveat:
        tf = box(s, 0.55, top, 12.25, 0.42)
        runs(tf, [("But:  ", RED, True), (caveat, GRAY, False)], size=12, first=True, space_after=0)
        top += 0.46
    fit_image(s, img, 0.35, top, 12.63, 7.2 - top)
    notes(s, note); return s


F = lambda n: os.path.join(FIGS, n)

# ---------------------------------------------------------------- 1 · title
s = slide()
tf = box(s, 0.9, 1.7, 11.5, 3.6)
para(tf, "Paris at Ten", size=46, bold=True, first=True, space_after=8)
para(tf, "Emissions growth has fallen by two thirds since the Paris Agreement. We ask what the global record can, and cannot, credit the treaty for.",
     size=19, color=GRAY, space_after=14)
runs(tf, [("182 countries · 99.8% of global CO2 · 1990-2021", GREEN, True)], size=14)
notes(s, "World-level paper. Country-level mechanism analysis is the follow-on paper, referenced once at the end.")

# ---------------------------------------------------------------- 2 · the question
s = slide()
kicker_title(s, "THE QUESTION", "Something changed after Paris. Was it Paris?")
tf = box(s, 0.55, 1.6, 7.4, 5.4)
para(tf, "The good news is real. Across essentially every country that reports, CO2 growth fell from +2.1% a year over the Rio-to-Paris quarter century to +0.7% a year since. Worst-case scenarios have moved off the table.",
     size=15, color=GRAY, first=True, space_after=14)
para(tf, "But the obvious reading — that the agreement did it — is exactly the reading the data is worst at supporting. Carbon intensity has been falling for fifty years. Renewable costs collapsed on almost the same schedule as the treaty. And a slowdown in emissions can come simply from a slowdown in the economy.",
     size=15, color=GRAY, space_after=14)
para(tf, "So we ask a narrow question: what actually changed after 2015, and is 2015 a special year in the record — or just a year that happens to sit near one?",
     size=15, color=INK, bold=True, space_after=0)
tf = box(s, 8.35, 1.6, 4.45, 5.4)
para(tf, "Why we look at intensity", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "Emissions have three moving parts:", size=13.5, color=GRAY, space_after=6)
para(tf, "economic growth  +  fuel mix  +  energy efficiency", size=13.5, color=INK, bold=True, space_after=8)
para(tf, "Growth pushes emissions up. The fuel mix and efficiency are the only two things that push them down — so a policy effect has to show up in one of them before it can show up in emissions at all.",
     size=13.5, color=GRAY, space_after=0)
notes(s, "Frame sympathetically, then set up the decomposition that structures the whole talk.")

# ---------------------------------------------------------------- 3 · what we did
s = slide()
kicker_title(s, "WHAT WE DID", "Three steps", kcolor=BLUE)
steps = [("1", "Build a consistent world series",
          "Take every country still reporting in 2021 with almost complete data, fill the handful of gaps, and add them up. That is 182 countries and 99.8% of global CO2. We stop at 2021 because most countries have not reported 2022-23 yet."),
         ("2", "Let the data pick the turning point",
          "Plot each series so a straight line means a steady percentage change per year, then fit two lines meeting at a kink. Try every possible kink year and keep the one that fits best. Nothing about a treaty enters this step."),
         ("3", "Ask whether the treaty year is special",
          "A gently curving trend makes almost any year look significant, so significance alone proves nothing. We rank 2015 against every other candidate year. A treaty earns credit only if its year beats the years around it.")]
x = 0.55
for n, t, d in steps:
    tf = box(s, x, 1.7, 3.95, 4.9)
    para(tf, n, size=34, bold=True, color=BLUE, first=True, space_after=6)
    para(tf, t, size=15.5, bold=True, space_after=8)
    para(tf, d, size=13, color=GRAY, space_after=0)
    x += 4.1
tf = box(s, 0.55, 6.55, 12.25, 0.6)
runs(tf, [("A turning point tells us when a trend changed — never, by itself, why. ", INK, True),
          ("That is the whole method.", GRAY, False)], size=13.5, first=True)
notes(s, "One methodology slide. If asked for detail: two-segment regression with an autocorrelation-robust test, BIC for model choice, and a placebo sweep over candidate years.")

# ---------------------------------------------------------------- 4 · the data
finding("WHAT WE DID · THE DATA", "We use every country that reports, for as long as it reports",
        "Near-complete coverage — and the last two years are dropped because most countries have not reported them.",
        F("w2_fig5_panel_construction.png"),
        note="182 countries, 99.8% of 2015 global CO2, 13 interpolated cells out of 5,824. The 2022-23 cliff is a reporting lag, not a real change, and including it would badly distort the post-Paris window.",
        kcolor=BLUE)

# ---------------------------------------------------------------- 5 · finding 1
finding("WHAT WE FIND · 1 OF 3", "Emissions growth fell — but mostly because growth fell",
        "Two thirds of the slowdown is slower economic growth. The fuel mix improved. Energy efficiency did not change at all.",
        F("w2_fig1_what_changed.png"),
        note="Before → after: GDP +3.1 → +2.1; fuel mix -0.0 → -0.4; efficiency -1.0 → -1.0; CO2 +2.1 → +0.7. The efficiency result is the surprise: fifty years of steady improvement, entirely unaffected by the treaty era.")

# ---------------------------------------------------------------- 6 · finding 2
finding("WHAT WE FIND · 2 OF 3", "The fuel mix broke away from its trend. Efficiency stayed on it.",
        "Cleaner energy is the one place the world is now beating its own long-run trend — and it is where climate policy would show up first.",
        F("w2_fig2_anniversary_test.png"),
        note="Fuel mix -0.7 pp/yr faster than trend (t = 3.4); overall intensity -1.0 (t = 2.3); efficiency -0.3 (t = 1.3, not significant). Supply-side change, not demand-side.")

# ---------------------------------------------------------------- 7 · finding 3
finding("WHAT WE FIND · 3 OF 3", "Rio fails the placebo test. Paris passes it.",
        "Paris is the first treaty date the global record does not reject — it ranks in the top tenth of all candidate years.",
        F("w2_fig3_rio_vs_paris.png"),
        caveat="the evidence is a plateau across 2012-2016, not a spike at 2015 — and that window coincides exactly with the collapse in clean-energy costs.",
        note="Rio: 14th percentile for overall intensity, 66th for fuel mix. Paris: 97th and 93rd. This is the contrast that makes the paper — the same test that dismisses Rio does not dismiss Paris.",
        kcolor=AMBER)

# ---------------------------------------------------------------- 9 · verdict
s = slide()
kicker_title(s, "WHAT IT ADDS UP TO", "Consistent with Paris. Not attributable to it.")
tf = box(s, 0.55, 1.6, 6.05, 5.4)
para(tf, "What we can say", size=15, bold=True, color=GREEN, first=True, space_after=6)
para(tf, "Emissions growth has fallen by two thirds, and part of that is a genuine improvement in the world's fuel mix that beats fifty years of trend. That improvement sits in the supply side, exactly where climate policy and clean technology operate. And 2015 ranks near the top of all candidate turning-point years — a test the 1992 Rio agreement fails badly.",
     size=13.5, color=GRAY, space_after=0)
tf = box(s, 7.0, 1.6, 5.8, 5.4)
para(tf, "What we cannot", size=15, bold=True, color=RED, first=True, space_after=6)
para(tf, "Most of the emissions slowdown is slower economic growth, not decarbonization. Energy efficiency did not respond at all. And the evidence for a turn is spread across 2012 to 2016 rather than concentrated on the treaty year — a window the test cannot resolve, and one that coincides exactly with the collapse in clean-energy costs.",
     size=13.5, color=GRAY, space_after=10)
para(tf, "The one thing that changed after Paris is the one thing cheap clean energy would also have changed.",
     size=13.5, color=INK, bold=True, space_after=0)
notes(s, "If one sentence survives, it is the bolded one. Neither vindication nor indictment.")

# ---------------------------------------------------------------- 10 · what would settle it
s = slide()
kicker_title(s, "WHAT WOULD SETTLE IT", "Three things — and two of them are not at the world level")
tf = box(s, 0.55, 1.65, 11.9, 5.0)
para(tf, "Persistence. The size of the change is already detectable — the test could pick up an effect half this large within a couple more years of data. What it cannot yet tell us is whether the new fuel-mix trend holds. If it survives to Paris+20, a one-off technology shock becomes a much weaker explanation.",
     size=15, color=GRAY, first=True, space_after=12)
para(tf, "The demand side. Efficiency is the more informative half precisely because it has not moved: cheap solar explains a cleaner fuel mix, but it does not explain how much energy the world uses per unit of output. If efficiency accelerates over the next decade, that is much harder to attribute to technology alone.",
     size=15, color=GRAY, space_after=12)
para(tf, "Disaggregation. Paris works through national pledges, so its signature should appear country by country — where series are longer, placebo tests are independent, and a documentary record exists to say what accompanied each turn. That is the companion paper.",
     size=15, color=GRAY, space_after=14)
para(tf, "The question we would put to the field: should treaty assessment be built on intensity turning points and placebo discipline, rather than on emissions levels and anniversaries?",
     size=15, color=INK, bold=True, space_after=0)
notes(s, "Forward-looking. Hands off to the national paper without previewing its results.")

# ---------------------------------------------------------------- 11 · limits
s = slide()
kicker_title(s, "LIMITS", "What we are not claiming", kcolor=RED)
col = [("This is not causal",
        "A turning point identifies when a trend changed, never why. Even a perfect fit at 2015 would be a coincidence of timing. And an agreement that prevented backsliding would leave no turning point at all — so the absence of one is not the absence of an effect."),
       ("The window matters",
        "Paris ranks 97th percentile against candidate years on a 1980-start series and lower on a 1990-start one; the fuel-mix result holds above the 90th in both. We report both rather than the flattering one."),
       ("Coverage and vintage",
        "182 countries and 99.8% of CO2, but only 13 gaps interpolated and the series stops in 2021 because 2022-23 reporting is incomplete. Adding those two years on partial coverage would overstate the acceleration."),
       ("Production-based accounting",
        "Emissions are counted where they are produced. At the world level that is the right frame — global production equals global consumption — but it matters for the national work that follows.")]
xs, ys = [0.55, 7.0, 0.55, 7.0], [1.6, 1.6, 4.35, 4.35]
for (t, d), x, y in zip(col, xs, ys):
    tf = box(s, x, y, 5.9, 2.6)
    para(tf, t, size=14.5, bold=True, first=True, space_after=4)
    para(tf, d, size=12, color=GRAY, space_after=0)
notes(s, "Volunteer the window sensitivity — a referee will look for it.")

# ---------------------------------------------------------------- 12 · close
s = slide()
tf = box(s, 0.9, 1.6, 11.5, 2.0)
para(tf, "The curve is bending.", size=28, bold=True, first=True, space_after=2)
para(tf, "The record cannot yet tell us who bent it.", size=28, bold=True, color=GREEN, space_after=0)
tf = box(s, 0.9, 3.7, 11.5, 2.9)
for t in ["Emissions growth fell by two thirds after Paris — but two thirds of that is slower economic growth.",
          "The fuel mix is the one component beating its long-run trend. Energy efficiency has not changed at all.",
          "Rio fails the placebo test; Paris passes it — the first treaty date the record does not reject.",
          "And the change coincides with the clean-energy cost collapse, which six years of data cannot separate from it."]:
    runs(tf, [("—  ", GREEN, True), (t, GRAY, False)], size=15, space_after=10)
tf = box(s, 0.9, 6.6, 11.5, 0.7)
para(tf, "Turning points say when · the Kaya split says through which door · placebo tests say whether the date is special",
     size=13, italic=True, color=LGRAY, first=True, align=PP_ALIGN.CENTER)
notes(s, "Close on the reframed headline.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
