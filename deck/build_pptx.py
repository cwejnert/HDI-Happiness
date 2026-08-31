"""
Render content.py as a 16:9 slide deck for presenting to the co-author team.

Same argument as build_artifact.py, paced for a meeting: act dividers, one
beat per slide, figures at full height beside the text.

    python build_pptx.py [out.pptx]

Default output: deck/The_Education_Exception.pptx
"""
from __future__ import annotations

import re
import sys
from html import unescape
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

import content as C

HERE = Path(__file__).parent
FIGDIR = HERE / "figures"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)

INK = RGBColor(0x14, 0x17, 0x1C)
INK_SOFT = RGBColor(0x5A, 0x60, 0x68)
INK_FAINT = RGBColor(0x86, 0x8D, 0x96)
ACCENT = RGBColor(0x17, 0x54, 0x9E)
ACCENT_PALE = RGBColor(0xC9, 0xD8, 0xEC)
FLAG = RGBColor(0xA6, 0x3A, 0x30)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
GROUND = RGBColor(0xE8, 0xEB, 0xEF)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"


def plain(s: str) -> str:
    """content.py carries HTML entities for the web build; strip them here."""
    return unescape(re.sub(r"<[^>]+>", "", s))


def fill(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size, color, font=SANS, bold=False, italic=False,
         space_after=0, space_before=0, spacing=None, caps=False, first=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    r = p.add_run()
    r.text = plain(text).upper() if caps else plain(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if spacing:
        p.line_spacing = spacing
    return p


def rule(slide, left, top, width, color=INK, thickness=Pt(1.5)):
    ln = slide.shapes.add_shape(1, left, top, width, thickness)  # rectangle
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------
def slide_title(prs):
    s = blank(prs)
    fill(s, GROUND)
    rule(s, MARGIN, Inches(1.5), W - 2 * MARGIN, ACCENT, Pt(2))

    tf = textbox(s, MARGIN, Inches(1.05), Inches(8), Inches(0.3))
    para(tf, "Commentary · working draft", size=11, color=ACCENT, font=MONO, first=True)

    tf = textbox(s, MARGIN, Inches(1.95), Inches(10), Inches(1.8))
    para(tf, C.TITLE, size=60, color=INK, font=SERIF, spacing=0.95, first=True)

    tf = textbox(s, MARGIN, Inches(3.85), Inches(9.2), Inches(1.3))
    para(tf, C.SUBTITLE, size=17, color=INK_SOFT, font=SERIF, italic=True,
         spacing=1.35, first=True)

    tf = textbox(s, MARGIN, Inches(6.1), Inches(11.6), Inches(1.0))
    para(tf, C.DATELINE, size=10.5, color=INK_SOFT, font=MONO, space_after=5, first=True)
    para(tf, C.SCOPE, size=9.5, color=INK_FAINT, font=MONO, spacing=1.3)


def slide_hook(prs):
    s = blank(prs)
    fill(s, SURFACE)
    tf = textbox(s, MARGIN, Inches(0.62), Inches(6.0), Inches(0.3))
    para(tf, C.HOOK["kicker"], size=11, color=INK_FAINT, font=MONO, caps=True, first=True)

    tf = textbox(s, MARGIN, Inches(1.12), Inches(6.0), Inches(1.7))
    para(tf, C.HOOK["heading"], size=30, color=INK, font=SERIF, spacing=1.06, first=True)

    tf = textbox(s, MARGIN, Inches(3.15), Inches(5.7), Inches(3.6))
    for i, p in enumerate(C.HOOK["body"]):
        para(tf, p, size=13, color=INK_SOFT, spacing=1.34, space_after=12, first=(i == 0))

    # the three acts, right column
    x = Inches(7.35)
    tf = textbox(s, x, Inches(0.62), Inches(5.3), Inches(0.3))
    para(tf, "The argument", size=11, color=ACCENT, font=MONO, caps=True, first=True)
    top = Inches(1.12)
    for numeral, claim, gloss in C.HOOK["acts"]:
        rule(s, x, top, Inches(5.3), RGBColor(0xCD, 0xD3, 0xDB), Pt(0.75))
        tf = textbox(s, x, top + Inches(0.16), Inches(0.55), Inches(0.5))
        para(tf, numeral, size=20, color=ACCENT, font=SERIF, first=True)
        tf = textbox(s, x + Inches(0.7), top + Inches(0.16), Inches(4.6), Inches(1.2))
        para(tf, claim, size=13.5, color=INK, bold=True, space_after=3, first=True)
        para(tf, gloss, size=11, color=INK_SOFT, spacing=1.25)
        top += Inches(1.45)


def slide_act_open(prs, act):
    s = blank(prs)
    fill(s, INK)

    tf = textbox(s, Inches(9.0), Inches(0.9), Inches(3.6), Inches(4.5))
    p = para(tf, act["numeral"], size=200, color=RGBColor(0x2A, 0x33, 0x40),
             font=SERIF, first=True)
    p.alignment = PP_ALIGN.RIGHT

    rule(s, MARGIN, Inches(1.5), Inches(8.0), ACCENT_PALE, Pt(2))

    tf = textbox(s, MARGIN, Inches(1.05), Inches(6.0), Inches(0.3))
    para(tf, f"Act {act['numeral']}", size=11, color=ACCENT_PALE, font=MONO,
         caps=True, first=True)

    tf = textbox(s, MARGIN, Inches(1.95), Inches(7.9), Inches(2.0))
    para(tf, act["title"], size=40, color=SURFACE, font=SERIF, spacing=1.02, first=True)

    tf = textbox(s, MARGIN, Inches(4.25), Inches(7.4), Inches(1.6))
    para(tf, act["thesis"], size=16, color=RGBColor(0x9A, 0xA4, 0xB0), font=SERIF,
         italic=True, spacing=1.35, first=True)

    # key numbers along the foot
    x = MARGIN
    colw = (W - 2 * MARGIN - Inches(0.6)) / 3
    for value, key in act["key_numbers"]:
        tf = textbox(s, x, Inches(6.15), colw, Inches(1.0))
        para(tf, value, size=22, color=SURFACE, font=MONO, space_after=4, first=True)
        para(tf, key, size=10, color=RGBColor(0x8A, 0x94, 0xA0), spacing=1.2)
        x += colw + Inches(0.3)


def slide_beat(prs, beat, numeral):
    s = blank(prs)
    fill(s, SURFACE)
    has_fig = bool(beat.get("figure"))
    textw = Inches(4.55) if has_fig else Inches(9.4)

    tf = textbox(s, MARGIN, Inches(0.55), textw, Inches(0.3))
    para(tf, beat["label"], size=10, color=FLAG, font=MONO, caps=True, first=True)

    tf = textbox(s, MARGIN, Inches(1.0), textw, Inches(1.5))
    para(tf, beat["heading"], size=23 if has_fig else 27, color=INK, font=SERIF,
         spacing=1.08, first=True)

    body_top = Inches(2.55) if has_fig else Inches(2.5)
    tf = textbox(s, MARGIN, body_top, textw, Inches(4.2))
    for i, p in enumerate(beat["body"]):
        para(tf, p, size=11.5 if has_fig else 14, color=INK_SOFT,
             spacing=1.32, space_after=11, first=(i == 0))

    tf = textbox(s, Inches(12.4), Inches(6.85), Inches(0.6), Inches(0.3))
    p = para(tf, numeral, size=11, color=INK_FAINT, font=MONO, first=True)
    p.alignment = PP_ALIGN.RIGHT

    if has_fig:
        img = FIGDIR / beat["figure"]
        from PIL import Image
        iw, ih = Image.open(img).size
        box_l, box_t = Inches(5.6), Inches(0.72)
        box_w, box_h = Inches(7.15), Inches(5.55)
        scale = min(box_w / iw, box_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        s.shapes.add_picture(str(img), box_l + int((box_w - w) / 2),
                             box_t + int((box_h - h) / 2), w, h)
        tf = textbox(s, box_l, Inches(6.45), box_w, Inches(0.8))
        para(tf, beat["caption"], size=9, color=INK_FAINT, spacing=1.25, first=True)


def slide_act_close(prs, act):
    s = blank(prs)
    fill(s, GROUND)
    rule(s, MARGIN, Inches(2.3), Inches(0.06), ACCENT, Inches(2.6))

    tf = textbox(s, Inches(1.0), Inches(1.75), Inches(9.5), Inches(0.3))
    para(tf, f"Where Act {act['numeral']} leaves us", size=11, color=ACCENT,
         font=MONO, caps=True, first=True)

    tf = textbox(s, Inches(1.0), Inches(2.35), Inches(10.4), Inches(3.2))
    para(tf, act["close"], size=21, color=INK, font=SERIF, spacing=1.35, first=True)


def slide_decisions(prs):
    s = blank(prs)
    fill(s, SURFACE)
    tf = textbox(s, MARGIN, Inches(0.55), Inches(9), Inches(0.9))
    para(tf, "Decisions for the team", size=32, color=INK, font=SERIF,
         space_after=6, first=True)
    para(tf, "Five things the acts do not settle, in the order they block drafting.",
         size=13, color=INK_SOFT)

    cols, x0, y0 = 3, MARGIN, Inches(2.0)
    colw = (W - 2 * MARGIN - Inches(0.7)) / cols
    for i, (head, body) in enumerate(C.DECISIONS):
        col, row = i % cols, i // cols
        x = x0 + col * (colw + Inches(0.35))
        y = y0 + row * Inches(2.45)
        rule(s, x, y, colw, ACCENT, Pt(2))
        tf = textbox(s, x, y + Inches(0.2), colw, Inches(2.0))
        para(tf, head, size=10.5, color=ACCENT, font=MONO, caps=True,
             space_after=7, first=True)
        para(tf, body, size=11, color=INK_SOFT, spacing=1.3)


def slide_appendix_divider(prs):
    s = blank(prs)
    fill(s, INK)
    tf = textbox(s, MARGIN, Inches(3.0), Inches(9), Inches(1.5))
    para(tf, "Evidence appendix", size=40, color=SURFACE, font=SERIF,
         space_after=10, first=True)
    para(tf, "Everything the three acts rest on but do not carry. "
             "In submission these become supplementary material.",
         size=14, color=RGBColor(0x9A, 0xA4, 0xB0), spacing=1.35)


def slide_appendix(prs, fname, title, desc):
    s = blank(prs)
    fill(s, SURFACE)
    tf = textbox(s, MARGIN, Inches(0.45), Inches(11.5), Inches(0.35))
    para(tf, "Appendix", size=10, color=INK_FAINT, font=MONO, caps=True, first=True)
    tf = textbox(s, MARGIN, Inches(0.82), Inches(11.5), Inches(0.5))
    para(tf, title, size=20, color=INK, font=SERIF, first=True)
    tf = textbox(s, MARGIN, Inches(1.32), Inches(11.5), Inches(0.5))
    para(tf, desc, size=11, color=INK_SOFT, spacing=1.28, first=True)

    from PIL import Image
    img = FIGDIR / fname
    iw, ih = Image.open(img).size
    box_l, box_t = MARGIN, Inches(2.15)
    box_w, box_h = W - 2 * MARGIN, Inches(4.9)
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(str(img), box_l + int((box_w - w) / 2), box_t, w, h)


def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else HERE / "The_Education_Exception.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_title(prs)
    slide_hook(prs)
    for act in C.ACTS:
        slide_act_open(prs, act)
        for beat in act["beats"]:
            slide_beat(prs, beat, act["numeral"])
        slide_act_close(prs, act)
    slide_decisions(prs)
    slide_appendix_divider(prs)
    for fname, title, desc in C.APPENDIX:
        slide_appendix(prs, fname, title, desc)

    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{out.stat().st_size / 1_048_576:.1f} MB)")


if __name__ == "__main__":
    main()
