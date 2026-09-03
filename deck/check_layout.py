"""Catch slide content that runs off the canvas.

PowerPoint text boxes do not clip, and this container cannot render a PPTX, so
an overlong body silently walks off the bottom edge with nothing to show for
it. build_pptx sizes bodies with fit_size, but a box whose ORIGIN plus declared
height already exceeds the slide is a layout bug no font shrinking can fix --
that is how a nine-item decisions grid at a fixed 2.45" row pitch escaped
review. This checks declared geometry for every shape, text boxes included.

    python check_layout.py [deck.pptx]
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

SLOP = Emu(9144)  # 0.01" -- rounding, not a real overrun


def main():
    path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "Measured_Where_It_Varies.pptx"
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            right, bottom = sh.left + (sh.width or 0), sh.top + (sh.height or 0)
            kind = "text" if sh.has_text_frame else "shape"
            label = (sh.text_frame.text[:40].replace("\n", " ") if sh.has_text_frame else "")
            if sh.left < -SLOP or sh.top < -SLOP:
                problems.append(f"slide {i:2d}  {kind} origin off-canvas "
                                f"({sh.left / 914400:.2f}, {sh.top / 914400:.2f})  {label}")
            if right > W + SLOP:
                problems.append(f"slide {i:2d}  {kind} runs off the right edge "
                                f"({right / 914400:.2f} > {W / 914400:.2f})  {label}")
            if bottom > H + SLOP:
                problems.append(f"slide {i:2d}  {kind} runs off the bottom "
                                f"({bottom / 914400:.2f} > {H / 914400:.2f})  {label}")
    print(f"{len(prs.slides)} slides checked against {W / 914400:.2f} x {H / 914400:.2f} in")
    if problems:
        print(f"\n{len(problems)} layout problems:")
        for p in problems:
            print(f"  {p}")
        return 1
    print("no shape exceeds the slide bounds")
    return 0


if __name__ == "__main__":
    sys.exit(main())
